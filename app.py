# Set up tiktoken encodings before any imports (PyInstaller compatibility)
import os
import sys
from pathlib import Path

def _get_tiktoken_path() -> Path:
    """Get tiktoken models directory path."""
    if getattr(sys, 'frozen', False) and hasattr(sys, '_MEIPASS'):
        return Path(sys._MEIPASS) / "models" / "tiktoken"
    return Path(__file__).parent / "models" / "tiktoken"

def _patch_tiktoken():
    """Patch tiktoken.load.read_file_cached to use local BPE files."""
    import tiktoken.load

    tiktoken_dir = _get_tiktoken_path()
    if not tiktoken_dir.exists():
        return

    _url_to_local = {
        "o200k_base.tiktoken": tiktoken_dir / "o200k_base.tiktoken",
        "cl100k_base.tiktoken": tiktoken_dir / "cl100k_base.tiktoken",
    }

    _original_read_file_cached = tiktoken.load.read_file_cached

    def _patched_read_file_cached(blobpath: str, expected_hash: str = None):
        for pattern, local_path in _url_to_local.items():
            if pattern in blobpath and local_path.exists():
                with open(local_path, "rb") as f:
                    return f.read()
        return _original_read_file_cached(blobpath, expected_hash)

    tiktoken.load.read_file_cached = _patched_read_file_cached

_patch_tiktoken()

import gradio as gr
import zipfile
import tempfile
import shutil
import json
from importlib import import_module
from llmWrapper.offline_translation import populate_sum_model
from typing import List, Tuple
from config.log_config import app_logger
import socket
import base64
import threading
import queue
import time
from functools import partial

# Import separated UI layout module
from ui_layout import (
    get_custom_css, create_header, create_footer, create_language_section,
    create_settings_section, create_model_glossary_section, create_main_interface,
    create_state_variables, create_theme_toggle, create_translation_history_button,
    create_history_page_content
)

# Import translation history manager
from config.translation_history import TranslationHistoryManager, format_duration, format_tokens

# Import language configs
from config.languages_config import LABEL_TRANSLATIONS, get_available_languages, get_language_code, add_custom_language

#-------------------------------------------------------------------------
# Constants and Configuration
#-------------------------------------------------------------------------

# File extension to translator module mapping
TRANSLATOR_MODULES = {
    ".docx": "translator.word_translator.WordTranslator",
    ".pptx": "translator.ppt_translator.PptTranslator",
    ".xlsx": "translator.excel_translator.ExcelTranslator",
    ".pdf": "translator.pdf_translator.PdfTranslator",
    ".srt": "translator.subtitle_translator.SubtitlesTranslator",
    ".txt": "translator.txt_translator.TxtTranslator",
    ".md": "translator.md_translator.MdTranslator",
}

# Note: Alternative translator modules have been merged into the main classes
# Excel: use_xlwings and bilingual_mode parameters
# Word: bilingual_mode parameter

# Global task queue and counter
task_queue = queue.Queue()
active_tasks = 0
task_lock = threading.Lock()

# Global variables for stop functionality
translation_stop_requested = False
current_translation_task = None
stop_lock = threading.Lock()

def generate_api_key_translations_js():
    """Generate JavaScript translations object from LABEL_TRANSLATIONS for API key section"""
    js_translations = {}
    for lang_code, translations in LABEL_TRANSLATIONS.items():
        js_translations[lang_code] = {
            "label": translations.get("API Key", "API Key"),
            "tooltipTitle": translations.get("Security Tips", "Security Tips"),
            "tooltipContent": translations.get("Security Tips Content", "")
        }
    return json.dumps(js_translations, ensure_ascii=False)

def enqueue_task(
    translate_func, files, model, src_lang, dst_lang,
    use_online, api_key, max_retries, max_token, thread_count, excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_name, session_lang, progress
):
    """Enqueue translation task or execute immediately if no tasks running"""
    global active_tasks
    
    with task_lock:
        if active_tasks == 0:
            # No active tasks, start immediately
            active_tasks += 1
            return None
        else:
            # Tasks running, add to queue
            task_info = {
                "files": files,
                "model": model,
                "src_lang": src_lang,
                "dst_lang": dst_lang,
                "use_online": use_online,
                "api_key": api_key,
                "max_retries": max_retries,
                "max_token": max_token,
                "thread_count": thread_count,
                "excel_mode_2": excel_mode_2,
                "excel_bilingual_mode": excel_bilingual_mode,
                "word_bilingual_mode": word_bilingual_mode,
                "pdf_bilingual_mode": pdf_bilingual_mode,
                "glossary_name": glossary_name,
                "session_lang": session_lang
            }
            task_queue.put(task_info)
            queue_position = task_queue.qsize()
            return f"Task added to queue. Position: {queue_position}"
        
def clean_server_cache():
    """In server_mode, clean temp/result/log dirs to prevent disk overflow.
    Called at the start of each translation to remove previous results."""
    if not server_mode:
        return
    try:
        temp_dir, result_dir, log_dir = get_custom_paths()
        for dir_path in [temp_dir, result_dir, log_dir]:
            if os.path.exists(dir_path):
                shutil.rmtree(dir_path, ignore_errors=True)
                os.makedirs(dir_path, exist_ok=True)
        # Also clean Gradio upload cache
        clean_gradio_cache()
        app_logger.info("Server cache cleaned")
    except Exception as e:
        app_logger.warning(f"Server cache cleanup error: {e}")

def clean_gradio_cache():
    """Clean up old Gradio temporary files"""
    try:
        gradio_temp_dir = tempfile.gettempdir()
        cleaned_count = 0
        
        for item in os.listdir(gradio_temp_dir):
            if item.startswith('gradio') or item.startswith('tmp'):
                item_path = os.path.join(gradio_temp_dir, item)
                try:
                    if time.time() - os.path.getmtime(item_path) > 300:
                        if os.path.isdir(item_path):
                            shutil.rmtree(item_path, ignore_errors=True)
                        else:
                            os.remove(item_path)
                        cleaned_count += 1
                except Exception as e:
                    app_logger.debug(f"Could not remove {item_path}: {e}")
        
        if cleaned_count > 0:
            app_logger.info(f"Cleaned {cleaned_count} Gradio cache items")
    except Exception as e:
        app_logger.warning(f"Gradio cache cleanup error: {e}")

def process_task_with_queue(
    translate_func, files, model, src_lang, dst_lang,
    use_online, api_key, max_retries, max_token, thread_count, excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_name, session_lang, progress
):
    """Process translation task and handle queue management"""
    global active_tasks
    if progress is None:
        progress = gr.Progress(track_tqdm=True)
    
    queue_msg = enqueue_task(
        translate_func, files, model, src_lang, dst_lang,
        use_online, api_key, max_retries, max_token, thread_count, excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_name, session_lang, progress
    )

    labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])
    stop_text = labels.get("Stop Translation", "Stop Translation")
    
    if queue_msg:
        return gr.update(value=None, visible=False), queue_msg, gr.update(value=stop_text, interactive=False)
    
    try:
        # Check if stop was requested before starting translation
        check_stop_requested()
        
        result = translate_func(
            files, model, src_lang, dst_lang,
            use_online, api_key, max_retries, max_token, thread_count, excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_name, session_lang, progress
        )
        process_next_task_in_queue(translate_func, progress)
        
        return result[0], result[1], result[2]
    except Exception as e:
        with task_lock:
            active_tasks -= 1
        process_next_task_in_queue(translate_func, progress)
        return gr.update(value=None, visible=False), f"Error: {str(e)}", gr.update(value=stop_text, interactive=False)

def process_next_task_in_queue(translate_func, progress):
    """Process next task in queue if available"""
    global active_tasks
    
    with task_lock:
        active_tasks -= 1
        
        if not task_queue.empty():
            next_task = task_queue.get()
            active_tasks += 1
            threading.Thread(
                target=process_queued_task,
                args=(translate_func, next_task, progress),
                daemon=True
            ).start()

def process_queued_task(translate_func, task_info, progress):
    """Process task from queue in separate thread"""
    try:
        # Check if stop was requested before starting
        check_stop_requested()
        
        if progress is None:
            progress = gr.Progress(track_tqdm=True)
        result = translate_func(
            task_info["files"],
            task_info["model"],
            task_info["src_lang"],
            task_info["dst_lang"],
            task_info["use_online"],
            task_info["api_key"],
            task_info["max_retries"],
            task_info["max_token"],
            task_info["thread_count"],
            task_info["excel_mode_2"],
            task_info["excel_bilingual_mode"],
            task_info["word_bilingual_mode"],
            task_info.get("pdf_bilingual_mode", False),
            task_info["glossary_name"],
            task_info.get("session_lang", "en"),
            progress
        )    
    except Exception as e:
        app_logger.exception(f"Error processing queued task: {e}")
    finally:
        process_next_task_in_queue(translate_func, progress)

class StopTranslationException(Exception):
    """Custom exception for when translation is stopped by user"""
    pass

def request_stop_translation(session_lang):
    """Request to stop current translation"""
    global translation_stop_requested
    
    labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])
    stopping_text = labels.get("Stopping", "Stopping...")
    
    with stop_lock:
        translation_stop_requested = True
    
    return gr.update(value=stopping_text, interactive=False)

def reset_stop_flag():
    """Reset stop flag for new translations"""
    global translation_stop_requested
    
    with stop_lock:
        translation_stop_requested = False

def check_stop_requested():
    """Check if stop has been requested"""
    with stop_lock:
        if translation_stop_requested:
            raise StopTranslationException("Translation stopped by user")
        return False

def modified_translate_button_click(
    translate_files_func, files, model, src_lang, dst_lang,
    use_online, api_key, max_retries, max_token, thread_count, excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_name,
    session_lang, continue_mode=False, progress=gr.Progress(track_tqdm=True)
):
    """Modified translate button click handler using task queue"""
    global current_translation_task

    labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])
    stop_text = labels.get("Stop Translation", "Stop Translation")

    # Reset UI and stop flag
    output_file_update = gr.update(visible=False)
    status_message = None
    reset_stop_flag()

    if not files:
        return output_file_update, "Please select file(s) to translate.", gr.update(value=stop_text, interactive=False)

    # In server_mode, API key comes from environment variable, skip client-side check
    if use_online and not api_key and not server_mode:
        return output_file_update, "API key is required for online models.", gr.update(value=stop_text, interactive=False)

    def wrapped_translate_func(files, model, src_lang, dst_lang,
                              use_online, api_key, max_retries, max_token, thread_count,
                              excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_name, session_lang, progress):
        return translate_files_func(files, model, src_lang, dst_lang,
                                   use_online, api_key, max_retries, max_token, thread_count,
                                   excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_name, session_lang,
                                   continue_mode=continue_mode, progress=progress)
    
    return process_task_with_queue(
        wrapped_translate_func, files, model, src_lang, dst_lang,
        use_online, api_key, max_retries, max_token, thread_count, excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_name, session_lang, progress
    )

def check_temp_translation_exists(files):
    """Check if temporary translation folders exist in custom temp directory"""
    if not files:
        return False, "No files selected."
    
    # Use custom temp directory from config
    temp_base_dir, _, _ = get_custom_paths()
    os.makedirs(temp_base_dir, exist_ok=True)
    
    found_folders = []
    
    for file_obj in files:
        # Get filename without extension
        filename = os.path.splitext(os.path.basename(file_obj.name))[0]
        
        # Look for exact matching folder in temp directory
        temp_folder = os.path.join(temp_base_dir, filename)
        
        if os.path.exists(temp_folder) and os.path.isdir(temp_folder):
            found_folders.append(temp_folder)
    
    if found_folders:
        return True, f"Found {len(found_folders)} existing translation folders."
    else:
        return False, "No existing translations found."

#-------------------------------------------------------------------------
# System Configuration Functions
#-------------------------------------------------------------------------

def read_system_config():
    """Read system configuration from config file"""
    config_path = os.path.join("config", "system_config.json")
    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return {
            "lan_mode": False,
            "default_online": False,
            "max_token": 768,
            "show_model_selection": True,
            "show_mode_switch": True,
            "show_lan_mode": True,
            "show_max_retries": True,
            "show_thread_count": True,
            "show_glossary": True,
            "excel_mode_2": False,
            "excel_bilingual_mode": False,
            "word_bilingual_mode": False,
            "pdf_bilingual_mode": False,
            "default_thread_count_online": 2,
            "default_thread_count_offline": 4,
            "default_src_lang": "English",
            "default_dst_lang": "English",
            "temp_dir": "temp",
            "result_dir": "result",
            "log_dir": "log"
        }

def write_system_config(config):
    """Write system configuration to config file"""
    config_path = os.path.join("config", "system_config.json")
    os.makedirs(os.path.dirname(config_path), exist_ok=True)
    with open(config_path, 'w', encoding='utf-8') as f:
        json.dump(config, f, indent=4, ensure_ascii=False)

def get_custom_paths():
    """Get custom directory paths from config and ensure they exist"""
    config = read_system_config()
    temp_dir = config.get("temp_dir", "temp")
    result_dir = config.get("result_dir", "result")
    log_dir = config.get("log_dir", "log")
    
    # Ensure directories exist
    os.makedirs(temp_dir, exist_ok=True)
    os.makedirs(result_dir, exist_ok=True)
    os.makedirs(log_dir, exist_ok=True)
    
    return temp_dir, result_dir, log_dir

def update_lan_mode(lan_mode):
    """Update system config with new LAN mode setting"""
    config = read_system_config()
    config["lan_mode"] = lan_mode
    write_system_config(config)
    return config["lan_mode"]

def update_online_mode(use_online):
    """Update system config with new online mode setting"""
    config = read_system_config()
    config["default_online"] = use_online
    write_system_config(config)
    return config["default_online"]

def update_max_retries(max_retries):
    """Update system config with new max retries setting"""
    config = read_system_config()
    config["max_retries"] = max_retries
    write_system_config(config)
    return max_retries

def update_thread_count(thread_count):
    """Update system config with new thread count setting"""
    config = read_system_config()
    # Update appropriate thread count based on current mode
    if config.get("default_online", False):
        config["default_thread_count_online"] = thread_count
    else:
        config["default_thread_count_offline"] = thread_count
    write_system_config(config)
    return thread_count

def update_excel_mode(excel_mode_2):
    """Update system config with new Excel mode setting"""
    config = read_system_config()
    config["excel_mode_2"] = excel_mode_2
    write_system_config(config)
    return excel_mode_2

def update_excel_bilingual_mode_with_auto_mode2(excel_bilingual_mode):
    """Update system config with new Excel bilingual mode setting and auto-enable mode 2"""
    config = read_system_config()
    config["excel_bilingual_mode"] = excel_bilingual_mode
    if excel_bilingual_mode:
        config["excel_mode_2"] = True
    
    write_system_config(config)

    return excel_bilingual_mode, config["excel_mode_2"]

def update_word_bilingual_mode(word_bilingual_mode):
    """Update system config with new Word bilingual mode setting"""
    config = read_system_config()
    config["word_bilingual_mode"] = word_bilingual_mode
    write_system_config(config)
    return word_bilingual_mode

def update_pdf_bilingual_mode(pdf_bilingual_mode):
    """Update system config with new PDF bilingual mode setting"""
    config = read_system_config()
    config["pdf_bilingual_mode"] = pdf_bilingual_mode
    write_system_config(config)
    return pdf_bilingual_mode

def update_language_preferences(src_lang=None, dst_lang=None):
    """Update system config with new language preferences"""
    config = read_system_config()
    
    if src_lang is not None:
        config["default_src_lang"] = src_lang
    if dst_lang is not None:
        config["default_dst_lang"] = dst_lang
        
    write_system_config(config)
    return config.get("default_src_lang"), config.get("default_dst_lang")

def get_default_languages():
    """Get default source and target languages from config"""
    config = read_system_config()
    default_src = config.get("default_src_lang", "English")
    default_dst = config.get("default_dst_lang", "English")
    return default_src, default_dst

def is_add_custom_option(value):
    """Check if the selected value is an 'Add Custom Language' option in any language"""
    # Get all translations of "Add Custom Language"
    add_custom_translations = set()
    for lang_code, labels in LABEL_TRANSLATIONS.items():
        if "Add Custom Language" in labels:
            add_custom_translations.add(labels["Add Custom Language"])
    # Also add the English default
    add_custom_translations.add("+ Add Custom...")
    add_custom_translations.add("+ Add Custom…")
    return value in add_custom_translations

def on_src_language_change(src_lang):
    """Handler for source language dropdown change"""
    if not is_add_custom_option(src_lang):
        update_language_preferences(src_lang=src_lang)

    # Return UI update for custom language row
    return gr.update(visible=is_add_custom_option(src_lang))

def on_dst_language_change(dst_lang):
    """Handler for target language dropdown change"""
    if not is_add_custom_option(dst_lang):
        update_language_preferences(dst_lang=dst_lang)

    # Return UI update for custom language row
    return gr.update(visible=is_add_custom_option(dst_lang))

def find_available_port(start_port=9980, max_attempts=20):
    """Find available port starting from start_port"""
    for port in range(start_port, start_port + max_attempts):
        try:
            with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
                s.bind(("127.0.0.1", port))
                return port
        except OSError:
            continue
    raise RuntimeError("No available port found.")

def resource_path(relative_path):
    """Get absolute path to resource, works for dev and PyInstaller"""
    try:
        # PyInstaller creates temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

def load_application_icon(config):
    """Load application icon using img_path from system_config.json"""
    # Get icon path from config
    img_path = config.get("img_path", "img/ico.ico")
    
    # Define MIME types for different image formats
    mime_types = {
        'ico': 'image/x-icon',
        'png': 'image/png',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'gif': 'image/gif',
        'svg': 'image/svg+xml'
    }
    
    # Paths to try in order
    icon_paths_to_try = []
    
    # Try absolute path if img_path is absolute
    if os.path.isabs(img_path):
        icon_paths_to_try.append(img_path)
    
    # Try from current directory
    if not os.path.isabs(img_path):
        icon_paths_to_try.append(img_path)
    
    # Try from PyInstaller _MEIPASS
    try:
        # PyInstaller creates temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
        # If img_path is not absolute, add it to _MEIPASS path
        if not os.path.isabs(img_path):
            meipass_path = os.path.join(base_path, img_path)
            icon_paths_to_try.append(meipass_path)
    except Exception:
        # Not running from PyInstaller bundle
        pass
    
    # Add default img/ico.ico as last resort
    default_icon = "img/ico.ico"
    if img_path != default_icon:
        # Try from current directory
        if default_icon not in icon_paths_to_try:
            icon_paths_to_try.append(default_icon)
        
        # Try from _MEIPASS
        try:
            base_path = sys._MEIPASS
            default_meipass_path = os.path.join(base_path, default_icon)
            if default_meipass_path not in icon_paths_to_try:
                icon_paths_to_try.append(default_meipass_path)
        except Exception:
            pass
    
    # Try each path in order
    for icon_path in icon_paths_to_try:
        try:
            if os.path.isfile(icon_path):
                image_type = icon_path.split('.')[-1].lower()
                mime_type = mime_types.get(image_type, 'image/png')
                
                app_logger.info(f"Loading icon from: {icon_path}")
                with open(icon_path, "rb") as f:
                    encoded_image = base64.b64encode(f.read()).decode("utf-8")
                return encoded_image, mime_type
        except Exception as e:
            app_logger.warning(f"Failed to load icon from {icon_path}: {e}")
            # Try next path
    
    # If all else fails, log error
    app_logger.error("Failed to load any icon, application will run without an icon")
    return None, None

#-------------------------------------------------------------------------
# Glossary Management Functions
#-------------------------------------------------------------------------

def get_glossary_files():
    """Get all CSV files from glossary directory"""
    glossary_dir = "glossary"
    
    # Ensure glossary directory exists
    os.makedirs(glossary_dir, exist_ok=True)
    # Get all CSV files
    try:
        csv_files = [f for f in os.listdir(glossary_dir) if f.endswith('.csv')]
        # Sort files with Default.csv first
        csv_files.sort(key=lambda x: (x != 'Default.csv', x.lower()))
        return [os.path.splitext(f)[0] for f in csv_files]  # Return without .csv extension
    except Exception as e:
        app_logger.warning(f"Error reading glossary directory: {e}")
        return ["Default"]

def update_glossary_selection(glossary_name):
    """Update system config with selected glossary"""
    config = read_system_config()
    config["default_glossary"] = glossary_name
    write_system_config(config)
    return glossary_name

def get_default_glossary():
    """Get default glossary from config"""
    config = read_system_config()
    return config.get("default_glossary", "Default")

def upload_glossary_file(file_obj, session_lang):
    """Handle glossary file upload - auto-uploads when file is selected"""
    if not file_obj:
        labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])
        return gr.update(), labels.get("No file selected", "No file selected."), gr.update(visible=True)

    glossary_dir = "glossary"
    os.makedirs(glossary_dir, exist_ok=True)

    try:
        # Get original filename
        original_name = os.path.basename(file_obj.name)

        # Check if it's CSV file
        if not original_name.lower().endswith('.csv'):
            labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])
            return gr.update(), labels.get("Only CSV files are allowed", "Only CSV files are allowed."), gr.update(visible=True)

        # Copy file to glossary directory
        dest_path = os.path.join(glossary_dir, original_name)

        # If file already exists, add number suffix
        counter = 1
        base_name, ext = os.path.splitext(original_name)
        while os.path.exists(dest_path):
            new_name = f"{base_name}_{counter}{ext}"
            dest_path = os.path.join(glossary_dir, new_name)
            counter += 1

        # Copy file
        import shutil
        shutil.copy2(file_obj.name, dest_path)

        # Update glossary choices
        labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])
        add_glossary_label = labels.get("Add Glossary", "+ Add Glossary...")
        updated_choices = get_glossary_files() + [add_glossary_label]
        new_glossary_name = os.path.splitext(os.path.basename(dest_path))[0]

        # Update config with new glossary
        update_glossary_selection(new_glossary_name)

        success_msg = labels.get("Glossary uploaded successfully", "Glossary uploaded successfully") + f": {new_glossary_name}"

        # Auto-select the new glossary and hide upload area
        return gr.update(choices=updated_choices, value=new_glossary_name), success_msg, gr.update(visible=False)

    except Exception as e:
        app_logger.exception(f"Error uploading glossary file: {e}")
        labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])
        error_msg = labels.get("Error uploading file", "Error uploading file") + f": {str(e)}"
        return gr.update(), error_msg, gr.update(visible=True)

def is_add_glossary_option(value):
    """Check if the selected value is an 'Add Glossary' option in any language"""
    add_glossary_translations = set()
    for lang_code, labels in LABEL_TRANSLATIONS.items():
        if "Add Glossary" in labels:
            add_glossary_translations.add(labels["Add Glossary"])
    # Also add English defaults
    add_glossary_translations.add("+ Add Glossary...")
    add_glossary_translations.add("+")
    return value in add_glossary_translations

def on_glossary_change(glossary_value, session_lang):
    """Handle glossary selection change"""
    if is_add_glossary_option(glossary_value):
        # Show file upload area
        return gr.update(visible=True)
    else:
        # Update config and hide upload controls
        if glossary_value:
            update_glossary_selection(glossary_value)
        return gr.update(visible=False)
    
#-------------------------------------------------------------------------
# Translation History Functions
#-------------------------------------------------------------------------

def load_translation_history(session_lang):
    """Load and render translation history as HTML"""
    _, _, log_dir = get_custom_paths()
    history_manager = TranslationHistoryManager(log_dir=log_dir)
    records = history_manager.get_all_records(limit=50)

    labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])

    if not records:
        return f"<div class='history-no-records'>{labels.get('No translation records', 'No translation records')}</div>"

    html_parts = []
    for record in records:
        # Status icon
        status = record.get("status", "unknown")
        if status == "success":
            status_icon = "✅"
            status_text = labels.get("Success", "Success")
        elif status == "failed":
            status_icon = "❌"
            status_text = labels.get("Failed", "Failed")
        elif status == "stopped":
            status_icon = "⏹️"
            status_text = labels.get("Stopped", "Stopped")
        else:
            status_icon = "❓"
            status_text = status

        # Format time
        start_time = record.get("start_time", "")
        if start_time:
            try:
                from datetime import datetime
                dt = datetime.fromisoformat(start_time)
                formatted_time = dt.strftime("%Y-%m-%d %H:%M")
            except:
                formatted_time = start_time[:16] if len(start_time) >= 16 else start_time
        else:
            formatted_time = "-"

        # Format duration
        duration_seconds = record.get("duration_seconds", 0)
        formatted_duration = format_duration(duration_seconds)

        # Format tokens
        total_tokens = record.get("total_tokens", 0)
        formatted_tokens = format_tokens(total_tokens)

        # Language info
        src_lang_display = record.get("src_lang_display", record.get("src_lang", ""))
        dst_lang_display = record.get("dst_lang_display", record.get("dst_lang", ""))

        # Model info
        model = record.get("model", "")
        use_online = record.get("use_online", False)
        mode_text = labels.get("Online", "Online") if use_online else labels.get("Offline", "Offline")

        # File info
        input_file = record.get("input_file", "")
        output_file_path = record.get("output_file_path", "")
        log_file_path = record.get("log_file_path", "")

        # Escape paths for JavaScript
        output_folder = os.path.dirname(output_file_path).replace("\\", "\\\\").replace("'", "\\'") if output_file_path else ""
        log_folder = os.path.dirname(log_file_path).replace("\\", "\\\\").replace("'", "\\'") if log_file_path else ""

        html = f"""
        <div class="history-record">
            <div class="history-record-header">
                <span class="history-record-filename">📄 {input_file}</span>
                <span class="history-record-status" title="{status_text}">{status_icon}</span>
            </div>
            <div class="history-record-info">
                <div class="history-record-info-item">
                    <span>🕐</span>
                    <span>{labels.get('Time', 'Time')}: {formatted_time}</span>
                </div>
                <div class="history-record-info-item">
                    <span>⏱️</span>
                    <span>{labels.get('Duration', 'Duration')}: {formatted_duration}</span>
                </div>
                <div class="history-record-info-item">
                    <span>🔢</span>
                    <span>{labels.get('Tokens', 'Tokens')}: {formatted_tokens}</span>
                </div>
                <div class="history-record-info-item">
                    <span>🌐</span>
                    <span>{src_lang_display} → {dst_lang_display}</span>
                </div>
                <div class="history-record-info-item">
                    <span>🤖</span>
                    <span>{model} ({mode_text})</span>
                </div>
            </div>
            <div class="history-record-actions">
                <button class="history-action-btn" onclick="openFolder('{output_folder}')" {'disabled' if not output_folder else ''}>
                    📂 {labels.get('Open Output Folder', 'Open Output Folder')}
                </button>
                <button class="history-action-btn" onclick="openFolder('{log_folder}')" {'disabled' if not log_folder else ''}>
                    📋 {labels.get('Open Log Folder', 'Open Log Folder')}
                </button>
            </div>
        </div>
        """
        html_parts.append(html)

    return "".join(html_parts)


def toggle_history_panel(is_visible):
    """Toggle history panel visibility"""
    return gr.update(visible=not is_visible)


def open_folder_path(path):
    """Open folder in file explorer"""
    import subprocess
    import platform

    if not path or not os.path.exists(path):
        return

    system = platform.system()
    try:
        if system == "Windows":
            subprocess.run(["explorer", path], check=False)
        elif system == "Darwin":  # macOS
            subprocess.run(["open", path], check=False)
        else:  # Linux
            subprocess.run(["xdg-open", path], check=False)
    except Exception as e:
        app_logger.error(f"Error opening folder: {e}")


#-------------------------------------------------------------------------
# Language and Localization Functions
#-------------------------------------------------------------------------

def parse_accept_language(accept_language: str) -> List[Tuple[str, float]]:
    """Parse Accept-Language into (language, q) pairs"""
    if not accept_language:
        return []
    
    languages = []
    for item in accept_language.split(','):
        item = item.strip()
        if not item:
            continue
        if ';q=' in item:
            lang, q = item.split(';q=')
            q = float(q)
        else:
            lang = item
            q = 1.0
        languages.append((lang, q))
    
    return sorted(languages, key=lambda x: x[1], reverse=True)

def get_user_lang(request: gr.Request) -> str:
    """Return top user language code that matches LANGUAGE_MAP"""
    try:
        # Handle different types of headers objects
        if hasattr(request.headers, 'get'):
            accept_lang = request.headers.get("accept-language", "").lower()
        elif hasattr(request.headers, 'accept-language'):
            accept_lang = getattr(request.headers, 'accept-language', "").lower()
        else:
            accept_lang = ""
    except (AttributeError, TypeError):
        accept_lang = ""
    
    parsed = parse_accept_language(accept_lang)
    
    if not parsed:
        return "en"
    
    highest_lang, _ = parsed[0]
    highest_lang = highest_lang.lower()

    if highest_lang.startswith("ja"):
        return "ja"
    elif highest_lang.startswith(("zh-tw", "zh-hk", "zh-hant")):
        return "zh-Hant"
    elif highest_lang.startswith(("zh-cn", "zh-hans", "zh")):
        return "zh"
    elif highest_lang.startswith("es"):
        return "es"
    elif highest_lang.startswith("fr"):
        return "fr"
    elif highest_lang.startswith("de"):
        return "de"
    elif highest_lang.startswith("it"):
        return "it"
    elif highest_lang.startswith("pt"):
        return "pt"
    elif highest_lang.startswith("ru"):
        return "ru"
    elif highest_lang.startswith("ko"):
        return "ko"
    elif highest_lang.startswith("th"):
        return "th"
    elif highest_lang.startswith("vi"):
        return "vi"
    elif highest_lang.startswith("en"):
        return "en"

    return "en"

def set_labels(session_lang: str):
    """Update UI labels according to chosen language"""
    labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])

    file_upload_label = "Upload Files"
    if "Upload Files" in labels:
        file_upload_label = labels["Upload Files"]
    elif "Upload File" in labels:
        file_upload_label = labels["Upload File"] + "s"

    # Update dropdown choices with translated "Add Custom Language" option
    custom_label = labels.get("Add Custom Language", "+ Add Custom...")
    new_choices = get_available_languages() + [custom_label]

    return {
        src_lang: gr.update(label=labels["Source Language"], choices=new_choices),
        dst_lang: gr.update(label=labels["Target Language"], choices=new_choices),
        use_online_model: gr.update(label=labels["Use Online Model"]),
        lan_mode_checkbox: gr.update(label=labels["Local Network Mode (Restart to Apply)"]),
        model_choice: gr.update(label=labels["Models"]),
        glossary_choice: gr.update(label=labels.get("Glossary", "Glossary")),
        max_retries_slider: gr.update(label=labels["Max Retries"]),
        thread_count_slider: gr.update(label=labels["Thread Count"]),
        api_key_input: gr.update(label=labels["API Key"], placeholder=labels.get("Enter your API key here", "Enter your API key here")),
        remember_key_checkbox: gr.update(label=labels.get("Remember Key", "Remember Key")),
        file_input: gr.update(label=file_upload_label),
        output_file: gr.update(label=labels["Download Translated File"]),
        status_message: gr.update(label=labels["Status Message"]),
        translate_button: gr.update(value=labels["Translate"]),
        continue_button: gr.update(value=labels["Continue Translation"]),
        excel_mode_checkbox: gr.update(label=labels.get("Excel Mode", "Excel Mode")),
        excel_bilingual_checkbox: gr.update(label=labels.get("Excel Bilingual", "Excel Bilingual")),
        word_bilingual_checkbox: gr.update(label=labels.get("Word Bilingual", "Word Bilingual")),
        pdf_bilingual_checkbox: gr.update(label=labels.get("PDF Bilingual", "PDF Bilingual")),
        stop_button: gr.update(value=labels.get("Stop Translation", "Stop Translation")),
        custom_lang_input: gr.update(
            label=labels.get("New Language Name", "New language name"),
            placeholder=labels.get("Language Name Placeholder", "e.g. Klingon")
        ),
        add_lang_button: gr.update(value=labels.get("Create Language", "Create")),
        history_nav_btn: gr.update(value=f"📋 {labels.get('Translation History', 'Translation History')}"),
        history_back_btn: gr.update(value=f"← {labels.get('Back', 'Back')}"),
        history_refresh_btn: gr.update(value=f"🔄 {labels.get('Refresh Records', 'Refresh')}"),
        history_title: gr.update(value=f"<h2 style='text-align: center; margin: 20px 0;'>{labels.get('Translation History', 'Translation History')}</h2>")
    }

#-------------------------------------------------------------------------
# UI and Model Functions
#-------------------------------------------------------------------------

def update_model_list_and_api_input(use_online):
    """Switch model options and show/hide API Key, update config"""
    # Update system config with new online mode
    update_online_mode(use_online)
    config = read_system_config()

    # Get appropriate thread count based on mode
    thread_count = config.get("default_thread_count_online", 2) if use_online else config.get("default_thread_count_offline", 4)

    # Get saved API key if remember is enabled
    remember_api_key = config.get("remember_api_key", False)

    if use_online:
        if default_online_model and default_online_model in online_models:
            default_online_value = default_online_model
        else:
            default_online_value = online_models[0] if online_models else None

        # Load API key for the selected model
        saved_api_key = load_api_key_for_model(default_online_value) if remember_api_key else ""

        return (
            gr.update(choices=online_models, value=default_online_value),
            gr.update(visible=not server_mode),
            gr.update(value=saved_api_key),
            gr.update(value=thread_count)
        )
    else:
        if default_local_model and default_local_model in local_models:
            default_local_value = default_local_model
        else:
            default_local_value = local_models[0] if local_models else None
        return (
            gr.update(choices=local_models, value=default_local_value),
            gr.update(visible=False),
            gr.update(value=""),
            gr.update(value=thread_count)
        )


def get_mykeys_dir():
    """Get the mykeys directory path, create if not exists"""
    mykeys_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "mykeys")
    os.makedirs(mykeys_dir, exist_ok=True)
    return mykeys_dir


def sanitize_model_name(model_name):
    """Sanitize model name to create a valid filename"""
    if not model_name:
        return "default"
    # Remove invalid filename characters and replace spaces
    invalid_chars = '<>:"/\\|?*'
    sanitized = model_name
    for char in invalid_chars:
        sanitized = sanitized.replace(char, '_')
    # Replace parentheses and spaces
    sanitized = sanitized.replace('(', '').replace(')', '').replace(' ', '_')
    return sanitized.strip('_') or "default"


def load_api_key_for_model(model_name):
    """Load API key for a specific model from mykeys folder"""
    mykeys_dir = get_mykeys_dir()
    key_file = os.path.join(mykeys_dir, f"{sanitize_model_name(model_name)}.json")

    try:
        if os.path.exists(key_file):
            with open(key_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return data.get("api_key", "")
    except (json.JSONDecodeError, IOError) as e:
        app_logger.warning(f"Failed to load API key for {model_name}: {e}")

    return ""


def save_api_key_for_model(model_name, api_key):
    """Save API key for a specific model to mykeys folder"""
    mykeys_dir = get_mykeys_dir()
    key_file = os.path.join(mykeys_dir, f"{sanitize_model_name(model_name)}.json")

    try:
        data = {
            "model": model_name,
            "api_key": api_key
        }
        with open(key_file, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        app_logger.info(f"API key saved for model: {model_name}")
    except IOError as e:
        app_logger.error(f"Failed to save API key for {model_name}: {e}")


def delete_api_key_for_model(model_name):
    """Delete API key file for a specific model"""
    mykeys_dir = get_mykeys_dir()
    key_file = os.path.join(mykeys_dir, f"{sanitize_model_name(model_name)}.json")

    try:
        if os.path.exists(key_file):
            os.remove(key_file)
            app_logger.info(f"API key deleted for model: {model_name}")
    except IOError as e:
        app_logger.error(f"Failed to delete API key for {model_name}: {e}")


def update_remember_api_key(remember, api_key, model_name, lan_mode):
    """Update remember API key setting and save/delete key based on toggle"""
    # Only allow saving in non-LAN mode
    if lan_mode:
        return False

    config = read_system_config()
    config["remember_api_key"] = remember
    write_system_config(config)

    if remember and api_key and model_name:
        save_api_key_for_model(model_name, api_key)
    elif not remember and model_name:
        delete_api_key_for_model(model_name)

    return remember


def save_api_key_on_change(api_key, remember, model_name):
    """Save API key when it changes if remember is enabled"""
    if remember and model_name:
        save_api_key_for_model(model_name, api_key)


def load_api_key_on_model_change(model_name, remember):
    """Load API key when model changes"""
    if remember and model_name:
        saved_key = load_api_key_for_model(model_name)
        return gr.update(value=saved_key)
    return gr.update(value="")

def refresh_models(use_online):
    """Refresh model list by re-scanning available models"""
    global local_models, online_models

    # Re-populate local models (force refresh to rescan)
    local_models = populate_sum_model(force_refresh=True) or []

    # Online models are typically static, but we can refresh them too
    online_models = [
        "gpt-4o-mini", "gpt-4o", "gpt-4-turbo", "gpt-4",
        "gpt-3.5-turbo", "gpt-3.5-turbo-16k",
        "deepseek-chat", "deepseek-reasoner",
        "claude-sonnet-4-20250514", "claude-3-7-sonnet-20250219",
        "gemini-2.0-flash", "gemini-2.0-flash-lite", "gemini-1.5-pro", "gemini-1.5-flash"
    ]

    app_logger.info(f"Models refreshed. Local: {len(local_models)}, Online: {len(online_models)}")

    if use_online:
        return gr.update(choices=online_models, value=online_models[0] if online_models else None)
    else:
        return gr.update(choices=local_models, value=local_models[0] if local_models else None)

def init_ui(request: gr.Request):
    """Set user language and update labels on page load"""
    user_lang = get_user_lang(request)
    config = read_system_config()
    
    lan_mode_state = config.get("lan_mode", False)
    default_online_state = config.get("default_online", False)
    max_token_state = config.get("max_token", MAX_TOKEN)
    excel_mode_2_state = config.get("excel_mode_2", False)
    excel_bilingual_mode_state = config.get("excel_bilingual_mode", False)
    word_bilingual_mode_state = config.get("word_bilingual_mode", False)
    pdf_bilingual_mode_state = config.get("pdf_bilingual_mode", False)
    # Always use default 4 for max retries
    max_retries_state = 4
    
    # Get thread count based on mode
    thread_count_state = config.get("default_thread_count_online", 2) if default_online_state else config.get("default_thread_count_offline", 4)
    
    # Get visibility settings
    show_max_retries = config.get("show_max_retries", True)
    show_thread_count = config.get("show_thread_count", True)
    show_glossary = config.get("show_glossary", True)
    
    # Get default glossary with Add Glossary option
    default_glossary = get_default_glossary()
    labels = LABEL_TRANSLATIONS.get(user_lang, LABEL_TRANSLATIONS["en"])
    add_glossary_label = labels.get("Add Glossary", "+ Add Glossary...")
    glossary_choices = get_glossary_files() + [add_glossary_label]
    
    # Update use_online_model checkbox based on default_online setting
    use_online_value = default_online_state
    
    # Update model choices based on online/offline mode
    if use_online_value:
        model_choices = online_models
        if default_online_model and default_online_model in online_models:
            model_value = default_online_model
        else:
            model_value = online_models[0] if online_models else None
    else:
        model_choices = local_models
        if default_local_model and default_local_model in local_models:
            model_value = default_local_model
        else:
            model_value = local_models[0] if local_models else None
    
    label_updates = set_labels(user_lang)

    # Add visibility updates for max_retries, thread_count, and glossary
    label_updates[max_retries_slider] = gr.update(label=LABEL_TRANSLATIONS.get(user_lang, LABEL_TRANSLATIONS["en"])["Max Retries"], visible=show_max_retries)
    label_updates[thread_count_slider] = gr.update(label=LABEL_TRANSLATIONS.get(user_lang, LABEL_TRANSLATIONS["en"])["Thread Count"], visible=show_thread_count)

    # Update remember_key_checkbox - disable in LAN mode, set value from config
    remember_api_key = config.get("remember_api_key", False) if not lan_mode_state else False
    # Load API key for the current model from mykeys folder
    saved_api_key = load_api_key_for_model(model_value) if (remember_api_key and use_online_value) else ""
    label_updates[remember_key_checkbox] = gr.update(
        label=labels.get("Remember Key", "Remember Key"),
        value=remember_api_key,
        interactive=not lan_mode_state
    )
    # Update api_key_input with saved value if remember is enabled
    label_updates[api_key_input] = gr.update(
        label=labels["API Key"],
        placeholder=labels.get("Enter your API key here", "Enter your API key here"),
        value=saved_api_key
    )
    
    # Prepare return values
    label_values = list(label_updates.values())
    
    # Return settings values and UI updates
    return [
        user_lang, 
        lan_mode_state, 
        default_online_state,
        max_token_state,
        max_retries_state,
        excel_mode_2_state,
        excel_bilingual_mode_state,
        word_bilingual_mode_state,
        pdf_bilingual_mode_state,
        thread_count_state,
        use_online_value,
        gr.update(choices=model_choices, value=model_value),  # model_choice update
        gr.update(choices=glossary_choices, value=default_glossary, visible=show_glossary),  # glossary_choice update with visibility
        gr.update(visible=False)  # glossary_upload_row (initially hidden)
    ] + label_values

def get_default_dropdown_value(saved_lang, dropdown_choices):
    """Get appropriate default value for language dropdowns"""
    if saved_lang in dropdown_choices:
        return saved_lang
    return saved_lang

def show_mode_checkbox(files):
    """Show mode checkboxes based on file types: Excel, Word, PDF"""
    if not files:
        return gr.update(visible=False), gr.update(visible=False), gr.update(visible=False), gr.update(visible=False)

    # Check if at least one Excel file is present
    excel_files = [f for f in files if os.path.splitext(f.name)[1].lower() == ".xlsx"]
    excel_visible = bool(excel_files)

    # Check if at least one Word file is present
    word_files = [f for f in files if os.path.splitext(f.name)[1].lower() == ".docx"]
    word_visible = bool(word_files)

    # Check if at least one PDF file is present
    pdf_files = [f for f in files if os.path.splitext(f.name)[1].lower() == ".pdf"]
    pdf_visible = bool(pdf_files)

    return gr.update(visible=excel_visible), gr.update(visible=excel_visible), gr.update(visible=word_visible), gr.update(visible=pdf_visible)

def update_continue_button(files):
    """Check if temp folders exist for uploaded files and update continue button state"""
    if not files:
        return gr.update(interactive=False)
    
    # If multiple files selected, disable continue button
    if isinstance(files, list) and len(files) > 1:
        return gr.update(interactive=False)
    
    # Check if single file is PDF
    single_file = files[0] if isinstance(files, list) else files
    file_extension = os.path.splitext(single_file.name)[1].lower()
    
    # Disable continue button for PDF files
    if file_extension == ".pdf":
        return gr.update(interactive=False)
    
    # Only check for temp folders if single non-PDF file selected
    has_temp, _ = check_temp_translation_exists(files)
    return gr.update(interactive=has_temp)

#-------------------------------------------------------------------------
# Translation Processing Functions
#-------------------------------------------------------------------------

def get_translator_class(file_extension, excel_mode_2=False, word_bilingual_mode=False, excel_bilingual_mode=False, pdf_bilingual_mode=False):
    """Dynamically import and return appropriate translator class for file extension"""
    module_path = TRANSLATOR_MODULES.get(file_extension.lower())

    if not module_path:
        return None

    try:
        # Split into module path and class name
        module_name, class_name = module_path.rsplit('.', 1)

        # Import module
        module = import_module(module_name)

        # Get class
        translator_class = getattr(module, class_name)

        # For Excel, Word, and PDF, return a partial class with mode parameters
        if file_extension.lower() == ".xlsx":
            # Excel: use_xlwings for mode_2 or bilingual, bilingual_mode for bilingual
            return partial(translator_class,
                          use_xlwings=excel_mode_2 or excel_bilingual_mode,
                          bilingual_mode=excel_bilingual_mode)
        elif file_extension.lower() == ".docx":
            # Word: bilingual_mode for bilingual
            return partial(translator_class, bilingual_mode=word_bilingual_mode)
        elif file_extension.lower() == ".pdf":
            # PDF: word_bilingual_mode for bilingual output (dual PDF)
            return partial(translator_class, word_bilingual_mode=pdf_bilingual_mode)

        return translator_class
    except (ImportError, AttributeError) as e:
        app_logger.exception(f"Error importing translator for {file_extension}: {e}")
        return None

CHAR_LIMIT = 100_000
GITHUB_URL = "https://github.com/YANG-Haruka/LinguaHaru"

def count_file_chars(file_path):
    """Count characters in a file based on its type"""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext in (".txt", ".md", ".srt"):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return len(f.read())
        elif ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            return sum(len(p.text) for p in doc.paragraphs)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            total = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        total += sum(len(p.text) for p in shape.text_frame.paragraphs)
            return total
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            total = 0
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            total += len(str(cell))
            wb.close()
            return total
        elif ext == ".pdf":
            # Estimate from file size: ~2 chars per byte for text-heavy PDFs
            return os.path.getsize(file_path) // 2
    except Exception as e:
        app_logger.warning(f"Error counting chars for {file_path}: {e}")
    # Fallback: estimate from file size
    return os.path.getsize(file_path)

def translate_files(
    files, model, src_lang, dst_lang, use_online, api_key, max_retries=4, max_token=768, thread_count=4,
    excel_mode_2=False, excel_bilingual_mode=False, word_bilingual_mode=False, pdf_bilingual_mode=False, glossary_name="Default", session_lang="en", continue_mode=False, progress=gr.Progress(track_tqdm=True)
):
    """Translate one or multiple files using chosen model"""
    reset_stop_flag()  # Reset stop flag at beginning
    clean_server_cache()  # In server_mode, clean all previous cache
    clean_gradio_cache()

    labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])
    stop_text = labels.get("Stop Translation", "Stop Translation")

    if not files:
        return gr.update(value=None, visible=False), "Please select file(s) to translate.", gr.update(value=stop_text, interactive=False)

    if use_online and not api_key and not server_mode:
        return gr.update(value=None, visible=False), "API key is required for online models.", gr.update(value=stop_text, interactive=False)

    # Character limit check in server_mode
    if server_mode:
        total_chars = sum(count_file_chars(f.name) for f in (files if isinstance(files, list) else [files]))
        if total_chars > CHAR_LIMIT:
            limit_msg = labels.get("Char Limit Exceeded",
                f"File exceeds the {CHAR_LIMIT:,} character limit. Please use the GitHub project or Release version for larger files.")
            gr.Warning(limit_msg)
            return gr.update(value=None, visible=False), f"{limit_msg}\n{GITHUB_URL}", gr.update(value=stop_text, interactive=False)

    src_lang_code = get_language_code(src_lang)
    dst_lang_code = get_language_code(dst_lang)
    
    # Convert glossary name to full path
    glossary_path = os.path.join("glossary", f"{glossary_name}.csv") if glossary_name else None

    # Common progress callback function
    def progress_callback(progress_value, desc=None):
        if check_stop_requested():
            raise StopTranslationException("Translation stopped by user")
        progress(progress_value, desc=desc)

    try:
        # Check if multiple files or single file
        if isinstance(files, list) and len(files) > 1:
            result = process_multiple_files(
                files, model, src_lang_code, dst_lang_code,
                use_online, api_key, max_token, max_retries, thread_count, excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_path, continue_mode, progress_callback, session_lang
            )
        else:
            # Handle single file case
            single_file = files[0] if isinstance(files, list) else files
            result = process_single_file(
                single_file, model, src_lang_code, dst_lang_code,
                use_online, api_key, max_token, max_retries, thread_count, excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_path, continue_mode, progress_callback, session_lang
            )
        
        return result[0], result[1], gr.update(value=stop_text, interactive=False)
        
    except StopTranslationException:
        return gr.update(value=None, visible=False), "Translation stopped by user.", gr.update(value=stop_text, interactive=False)
    except Exception as e:
        return gr.update(value=None, visible=False), f"Error: {str(e)}", gr.update(value=stop_text, interactive=False)

def process_single_file(
    file, model, src_lang_code, dst_lang_code,
    use_online, api_key, max_token, max_retries, thread_count, excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_path, continue_mode, progress_callback, session_lang="en"
):
    """Process single file for translation"""
    file_name = os.path.basename(file.name)
    
    # Get custom paths from config
    temp_dir, result_dir, log_dir = get_custom_paths()
    
    # Create new log file for this file
    from config.log_config import file_logger
    file_logger.create_file_log(file_name, log_dir=log_dir)
    
    app_logger.info(f"Processing file: {file_name}")
    app_logger.info(f"Source language: {src_lang_code}, Target language: {dst_lang_code}, Model: {model}")
    
    file_name, file_extension = os.path.splitext(file.name)
    
    translator_class = get_translator_class(file_extension, excel_mode_2, word_bilingual_mode, excel_bilingual_mode, pdf_bilingual_mode)

    if not translator_class:
        return (
            gr.update(value=None, visible=False),
            f"Unsupported file type '{file_extension}'."
        )

    try:
        # Pass check_stop_requested function to translator with custom paths
        translator = translator_class(
            file.name, model, use_online, api_key,
            src_lang_code, dst_lang_code, continue_mode,
            max_token=max_token, max_retries=max_retries,
            thread_count=thread_count, glossary_path=glossary_path,
            temp_dir=temp_dir,      # Pass custom temp directory
            result_dir=result_dir,  # Pass custom result directory
            session_lang=session_lang,  # Pass session language for i18n
            log_dir=log_dir         # Pass custom log directory
        )
        
        # Add check_stop_requested as attribute
        translator.check_stop_requested = check_stop_requested

        # Get translated labels
        labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])

        progress_callback(0, desc=f"{labels.get('Extracting text', 'Extracting text')}...")

        translated_file_path, missing_counts = translator.process(
            file_name, file_extension, progress_callback=progress_callback
        )

        # Format completion message with tokens
        completion_msg = labels.get("Translation completed", "Translation completed")
        tokens_msg = labels.get("Total tokens used", "Total tokens used")

        # Get total tokens from translator
        total_tokens = getattr(translator, 'total_tokens', 0)
        if total_tokens > 0:
            # Format tokens with K suffix for thousands
            if total_tokens >= 1000:
                tokens_str = f"{total_tokens / 1000:.1f}K"
            else:
                tokens_str = str(total_tokens)
            final_msg = f"{completion_msg} | {tokens_msg}: {tokens_str}"
        else:
            final_msg = completion_msg

        progress_callback(1, desc=final_msg)

        if missing_counts:
            msg = f"Warning: Missing segments for keys: {sorted(missing_counts)}"
            return gr.update(value=translated_file_path, visible=True), msg

        return gr.update(value=translated_file_path, visible=True), final_msg
    
    except StopTranslationException:
        app_logger.info("Translation stopped by user")
        # Save stopped status to history
        if 'translator' in locals() and translator:
            translator.save_stopped_summary()
        return gr.update(value=None, visible=False), "Translation stopped by user."
    except ValueError as e:
        # Save failed status to history
        if 'translator' in locals() and translator:
            translator.save_failed_summary()
        return gr.update(value=None, visible=False), f"Translation failed: {str(e)}"
    except Exception as e:
        app_logger.exception("Error processing file")
        # Save failed status to history
        if 'translator' in locals() and translator:
            translator.save_failed_summary()
        return gr.update(value=None, visible=False), f"Error: {str(e)}"
    
def process_multiple_files(
    files, model, src_lang_code, dst_lang_code,
    use_online, api_key, max_token, max_retries, thread_count, excel_mode_2, excel_bilingual_mode, word_bilingual_mode, pdf_bilingual_mode, glossary_path, continue_mode, progress_callback, session_lang="en"
):
    """Process multiple files and return zip archive"""
    # Get custom paths from config
    temp_dir, result_dir, log_dir = get_custom_paths()
    
    # Create temporary directory for translated files in custom result directory
    temp_zip_dir = tempfile.mkdtemp(prefix="translated_", dir=result_dir)
    zip_path = os.path.join(temp_zip_dir, "translated_files.zip")
    
    try:
        valid_files = []
        
        # Validate all files
        for file_obj in files:
            _, ext = os.path.splitext(file_obj.name)
            if get_translator_class(ext, excel_mode_2, word_bilingual_mode, excel_bilingual_mode, pdf_bilingual_mode):
                file_name = os.path.basename(file_obj.name)
                valid_files.append((file_obj, file_name))
        
        if not valid_files:
            shutil.rmtree(temp_zip_dir)
            return gr.update(value=None, visible=False), "No supported files found."
        
        # Create zip file
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            total_files = len(valid_files)
            total_tokens = 0  # Track total tokens across all files
            
            for i, (file_obj, rel_path) in enumerate(valid_files):
                # Create new log file for current file being processed
                from config.log_config import file_logger
                file_logger.create_file_log(rel_path, log_dir=log_dir)
                
                app_logger.info(f"Processing file {i+1}/{total_files}: {rel_path}")
                
                file_name, file_extension = os.path.splitext(file_obj.name)
                base_name = os.path.basename(file_name)
                
                # Update progress with initial file info
                progress_callback(i / total_files, desc=f"Starting to process {rel_path} (File {i+1}/{total_files})")
                
                # Create translator for this file, passing mode parameters
                translator_class = get_translator_class(file_extension, excel_mode_2, word_bilingual_mode, excel_bilingual_mode, pdf_bilingual_mode)
                if not translator_class:
                    continue  # Skip unsupported files (should not happen due to earlier validation)
                
                try:
                    # Process file with custom paths
                    translator = translator_class(
                        file_obj.name, model, use_online, api_key,
                        src_lang_code, dst_lang_code, continue_mode, max_token=max_token, max_retries=max_retries,
                        thread_count=thread_count, glossary_path=glossary_path,
                        temp_dir=temp_dir,      # Pass custom temp directory
                        result_dir=result_dir,  # Pass custom result directory
                        session_lang=session_lang,  # Pass session language for i18n
                        log_dir=log_dir         # Pass custom log directory
                    )
                    
                    # Create output directory
                    output_dir = os.path.join(temp_zip_dir, "files")
                    os.makedirs(output_dir, exist_ok=True)
                    
                    # Create progress callback that shows individual file progress and overall position
                    def file_progress(value, desc=None):
                        file_desc = desc if desc else ""
                        overall_info = f" (File {i+1}/{total_files})"
                        progress_callback(i / total_files + value / total_files, desc=f"{file_desc}{overall_info}")
                    
                    translated_file_path, _ = translator.process(
                        os.path.join(output_dir, base_name),
                        file_extension,
                        progress_callback=file_progress
                    )
                    
                    # Add to zip
                    zipf.write(
                        translated_file_path, 
                        os.path.basename(translated_file_path)
                    )
                    # Accumulate total tokens
                    total_tokens += getattr(translator, 'total_tokens', 0)

                except StopTranslationException:
                    app_logger.info(f"Translation stopped by user for file {rel_path}")
                    if 'translator' in locals() and translator:
                        translator.save_stopped_summary()
                    # Re-raise to stop processing all files
                    raise
                except Exception as e:
                    app_logger.exception(f"Error processing file {rel_path}: {e}")
                    # Save failed status to history
                    if 'translator' in locals() and translator:
                        translator.save_failed_summary()
                    # Continue with next file

        # Get translated labels
        labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])
        completion_msg = labels.get("Translation completed", "Translation completed")
        tokens_label = labels.get("Total tokens used", "Total tokens used")

        # Format tokens
        if total_tokens > 0:
            if total_tokens >= 1000:
                tokens_str = f"{total_tokens / 1000:.1f}K"
            else:
                tokens_str = str(total_tokens)
            final_msg = f"{completion_msg} ({total_files} files) | {tokens_label}: {tokens_str}"
        else:
            final_msg = f"{completion_msg} ({total_files} files)"

        progress_callback(1, desc=final_msg)
        return gr.update(value=zip_path, visible=True), final_msg
    
    except Exception as e:
        app_logger.exception("Error processing files")
        shutil.rmtree(temp_zip_dir)
        return gr.update(value=None, visible=False), f"Error processing files: {str(e)}"

#-------------------------------------------------------------------------
# Main Application Initialization
#-------------------------------------------------------------------------

# Prevent re-initialization in subprocess (Windows multiprocessing spawns new processes)
import multiprocessing
_is_main_process = multiprocessing.current_process().name == 'MainProcess'

# Read initial configuration (needed for both main and subprocess)
config = read_system_config()
server_mode = config.get("server_mode", False)
initial_lan_mode = config.get("lan_mode", False)
initial_default_online = config.get("default_online", False)
initial_max_token = config.get("max_token", 768)
initial_max_retries = config.get("max_retries", 4)
initial_excel_mode_2 = config.get("excel_mode_2", False)
initial_excel_bilingual_mode = config.get("excel_bilingual_mode", False)
initial_word_bilingual_mode = config.get("word_bilingual_mode", False)
initial_thread_count_online = config.get("default_thread_count_online", 2)
initial_thread_count_offline = config.get("default_thread_count_offline", 4)
initial_thread_count = initial_thread_count_online if initial_default_online else initial_thread_count_offline
app_title = config.get("app_title", "LinguaHaru")
app_title_web = "LinguaHaru" if app_title == "" else app_title
img_path = config.get("img_path", "img/ico.png")
img_height = config.get("img_height", 250)

# Update global MAX_TOKEN from config
MAX_TOKEN = initial_max_token

# Get visibility settings from config
initial_show_model_selection = config.get("show_model_selection", True)
initial_show_mode_switch = config.get("show_mode_switch", True)
initial_show_lan_mode = config.get("show_lan_mode", True)
initial_show_max_retries = config.get("show_max_retries", True)
initial_show_thread_count = config.get("show_thread_count", True)
initial_show_glossary = config.get("show_glossary", True)
default_local_model = config.get("default_local_model", "")
default_online_model = config.get("default_online_model", "")

# Only run heavy initialization in main process
if _is_main_process:
    # Load local and online models
    local_models = populate_sum_model() or []
    CUSTOM_LABEL = "+ Add Custom…"
    dropdown_choices = get_available_languages() + [CUSTOM_LABEL]
    config_dir = "config/api_config"
    online_models = [
        os.path.splitext(f)[0] for f in os.listdir(config_dir)
        if f.endswith(".json") and f != "Custom.json"
    ]

    encoded_image, mime_type = load_application_icon(config)

    # Initialize custom directories
    get_custom_paths()
else:
    # Subprocess: use minimal initialization
    local_models = []
    online_models = []
    CUSTOM_LABEL = "+ Add Custom…"
    dropdown_choices = get_available_languages() + [CUSTOM_LABEL]
    encoded_image, mime_type = None, None

#-------------------------------------------------------------------------
# Gradio UI Construction
#-------------------------------------------------------------------------
# Create Gradio blocks interface with enhanced language dropdown styling
with gr.Blocks(
    title=app_title_web,
    css=get_custom_css()
) as demo:
    
    # Create theme toggle button first (positioned absolutely)
    theme_toggle_btn = create_theme_toggle()
    
    # Create header
    create_header(app_title, encoded_image, mime_type, img_height)
    
    # Create footer
    create_footer()
    
    # Create state variables
    states = create_state_variables(config)
    session_lang = states['session_lang']
    lan_mode_state = states['lan_mode_state']
    default_online_state = states['default_online_state']
    max_token_state = states['max_token_state']
    max_retries_state = states['max_retries_state']
    excel_mode_2_state = states['excel_mode_2_state']
    excel_bilingual_mode_state = states['excel_bilingual_mode_state']
    word_bilingual_mode_state = states['word_bilingual_mode_state']
    pdf_bilingual_mode_state = states['pdf_bilingual_mode_state']
    thread_count_state = states['thread_count_state']

    default_src_lang, default_dst_lang = get_default_languages()

    # Create get_label function for i18n (uses English for initial render, updated on page load)
    initial_labels = LABEL_TRANSLATIONS["en"]
    def get_label(key):
        return initial_labels.get(key, key)

    # Main page content wrapped in a Column for page navigation
    with gr.Column(visible=True, elem_id="main-page") as main_page:
        # Create language selection section
        src_lang, swap_button, dst_lang, custom_lang_input, add_lang_button, custom_lang_row = create_language_section(
            default_src_lang, default_dst_lang, get_label
        )

        # Create settings section
        (use_online_model, lan_mode_checkbox, max_retries_slider,
        thread_count_slider, excel_mode_checkbox, excel_bilingual_checkbox, word_bilingual_checkbox, pdf_bilingual_checkbox) = create_settings_section(config)

        # Create model and glossary section
        (model_choice, model_refresh_btn, glossary_choice, glossary_upload_row,
         glossary_upload_file) = create_model_glossary_section(
            config, local_models, online_models, get_glossary_files, get_default_glossary, get_label
        )

        # Create main interface
        (api_key_input, api_key_row, remember_key_checkbox, file_input, output_file, status_message,
         translate_button, continue_button, stop_button) = create_main_interface(config, get_label)

        # Create translation history navigation button (hidden in server_mode)
        history_nav_btn = create_translation_history_button(get_label)
        if server_mode:
            history_nav_btn.visible = False

    # Create history page (initially hidden)
    with gr.Column(visible=False, elem_id="history-page") as history_page:
        history_back_btn, history_refresh_btn, history_title, history_list = create_history_page_content(get_label)

    # Hidden components for folder opening functionality
    folder_path_input = gr.Textbox(visible=False, elem_id="folder-path-input")
    folder_open_trigger = gr.Button(visible=False, elem_id="folder-open-trigger")

    # Event handlers
    use_online_model.change(
        update_model_list_and_api_input,
        inputs=use_online_model,
        outputs=[model_choice, api_key_row, api_key_input, thread_count_slider]
    ).then(
        fn=None,
        inputs=None,
        outputs=None,
        js="""
        () => {
            console.log('Online model toggled, initializing features...');
            function initFeatures() {
                if (window.initApiKeyEyeToggle) {
                    window.initApiKeyEyeToggle();
                }
                if (window.updateApiKeyLanguage && window.currentApiKeyLang) {
                    window.updateApiKeyLanguage(window.currentApiKeyLang);
                }
                // Setup tooltip
                const helpWrapper = document.getElementById('api-help-wrapper');
                const tooltip = document.getElementById('api-tooltip');
                if (helpWrapper && tooltip && !helpWrapper._tooltipInitialized) {
                    helpWrapper._tooltipInitialized = true;
                    const showTooltip = () => {
                        const rect = helpWrapper.getBoundingClientRect();
                        const tooltipWidth = 280;
                        tooltip.style.visibility = 'hidden';
                        tooltip.style.opacity = '0';
                        tooltip.style.display = 'block';
                        const tooltipHeight = tooltip.offsetHeight;
                        tooltip.style.display = '';
                        const gap = 8;
                        const iconCenterX = rect.left + (rect.width / 2);
                        let left = iconCenterX - (tooltipWidth / 2);
                        let top = rect.top - tooltipHeight - gap;
                        let arrowLeft = 50;
                        if (left < 10) {
                            arrowLeft = ((iconCenterX - 10) / tooltipWidth) * 100;
                            left = 10;
                        }
                        if (left + tooltipWidth > window.innerWidth - 10) {
                            const newLeft = window.innerWidth - tooltipWidth - 10;
                            arrowLeft = ((iconCenterX - newLeft) / tooltipWidth) * 100;
                            left = newLeft;
                        }
                        arrowLeft = Math.max(15, Math.min(85, arrowLeft));
                        if (top < 10) top = rect.bottom + gap;
                        tooltip.style.left = left + 'px';
                        tooltip.style.top = top + 'px';
                        tooltip.style.setProperty('--arrow-left', arrowLeft + '%');
                        tooltip.classList.add('visible');
                    };
                    const hideTooltip = () => {
                        tooltip.classList.remove('visible');
                    };
                    helpWrapper.addEventListener('mouseenter', showTooltip);
                    helpWrapper.addEventListener('mouseleave', hideTooltip);
                }
            }
            setTimeout(initFeatures, 300);
            setTimeout(initFeatures, 800);
        }
        """
    )

    # Remember API key checkbox handler
    remember_key_checkbox.change(
        update_remember_api_key,
        inputs=[remember_key_checkbox, api_key_input, model_choice, lan_mode_state],
        outputs=remember_key_checkbox
    )

    # Save API key when it changes (if remember is enabled)
    api_key_input.change(
        save_api_key_on_change,
        inputs=[api_key_input, remember_key_checkbox, model_choice],
        outputs=None
    )

    # Load API key when model changes (if remember is enabled)
    model_choice.change(
        load_api_key_on_model_change,
        inputs=[model_choice, remember_key_checkbox],
        outputs=api_key_input
    )

    # Model refresh button
    model_refresh_btn.click(
        refresh_models,
        inputs=use_online_model,
        outputs=model_choice
    )

    # Add LAN mode - also disable remember key when in LAN mode
    def update_lan_mode_with_remember_key(lan_mode):
        lan_state = update_lan_mode(lan_mode)
        # Disable remember key checkbox when in LAN mode (security)
        # Note: We don't delete the saved keys in mykeys folder, just disable the feature
        if lan_mode:
            config = read_system_config()
            config["remember_api_key"] = False
            write_system_config(config)
            return lan_state, gr.update(value=False, interactive=False)
        else:
            return lan_state, gr.update(interactive=True)

    lan_mode_checkbox.change(
        update_lan_mode_with_remember_key,
        inputs=lan_mode_checkbox,
        outputs=[lan_mode_state, remember_key_checkbox]
    )
    
    # Add Max Retries
    max_retries_slider.change(
        update_max_retries,
        inputs=max_retries_slider,
        outputs=max_retries_state
    )
    
    # Add Thread Count
    thread_count_slider.change(
        update_thread_count,
        inputs=thread_count_slider,
        outputs=thread_count_state
    )

    excel_mode_checkbox.change(
        update_excel_mode,
        inputs=excel_mode_checkbox,
        outputs=excel_mode_2_state
    )

    excel_bilingual_checkbox.change(
        update_excel_bilingual_mode_with_auto_mode2,
        inputs=excel_bilingual_checkbox,
        outputs=[excel_bilingual_mode_state, excel_mode_2_state]
    ).then(
        lambda excel_bilingual_mode: gr.update(value=True) if excel_bilingual_mode else gr.update(),
        inputs=excel_bilingual_checkbox,
        outputs=excel_mode_checkbox
    )

    word_bilingual_checkbox.change(
        update_word_bilingual_mode,
        inputs=word_bilingual_checkbox,
        outputs=word_bilingual_mode_state
    )

    pdf_bilingual_checkbox.change(
        update_pdf_bilingual_mode,
        inputs=pdf_bilingual_checkbox,
        outputs=pdf_bilingual_mode_state
    )

    file_input.change(
        fn=lambda files: [show_mode_checkbox(files)[0],
                        show_mode_checkbox(files)[1],
                        show_mode_checkbox(files)[2],
                        show_mode_checkbox(files)[3],
                        update_continue_button(files)],
        inputs=file_input,
        outputs=[excel_mode_checkbox, excel_bilingual_checkbox, word_bilingual_checkbox, pdf_bilingual_checkbox, continue_button]
    )

    # Glossary event handlers (only if glossary visible)
    if initial_show_glossary:
        glossary_choice.change(
            on_glossary_change,
            inputs=[glossary_choice, session_lang],
            outputs=[glossary_upload_row]
        )

        # Auto-upload when file is selected (no button click needed)
        glossary_upload_file.change(
            upload_glossary_file,
            inputs=[glossary_upload_file, session_lang],
            outputs=[glossary_choice, status_message, glossary_upload_row]
        )

    # Update event handlers for translate button
    translate_button.click(
        lambda: (gr.update(visible=False), None, gr.update(interactive=False), gr.update(interactive=False), gr.update(interactive=True)),
        inputs=[],
        outputs=[output_file, status_message, translate_button, continue_button, stop_button]
    ).then(
        partial(modified_translate_button_click, translate_files),
        inputs=[
            file_input, model_choice, src_lang, dst_lang,
            use_online_model, api_key_input, max_retries_slider, max_token_state,
            thread_count_slider, excel_mode_checkbox, excel_bilingual_checkbox, word_bilingual_checkbox, pdf_bilingual_checkbox, glossary_choice, session_lang
        ],
        outputs=[output_file, status_message, stop_button]
    ).then(
        lambda session_lang: (
            gr.update(interactive=True), 
            gr.update(interactive=True), 
            gr.update(value=LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"]).get("Stop Translation", "Stop Translation"), interactive=False)
        ),
        inputs=[session_lang],
        outputs=[translate_button, continue_button, stop_button]
    )

    # In continue_button.click event:
    continue_button.click(
        lambda: (gr.update(visible=False), None, gr.update(interactive=False), gr.update(interactive=False), gr.update(interactive=True)),
        inputs=[],
        outputs=[output_file, status_message, translate_button, continue_button, stop_button]
    ).then(
        partial(modified_translate_button_click, translate_files, continue_mode=True),
        inputs=[
            file_input, model_choice, src_lang, dst_lang,
            use_online_model, api_key_input, max_retries_slider, max_token_state,
            thread_count_slider, excel_mode_checkbox, excel_bilingual_checkbox, word_bilingual_checkbox, pdf_bilingual_checkbox, glossary_choice, session_lang
        ],
        outputs=[output_file, status_message, stop_button]
    ).then(
        lambda session_lang: (
            gr.update(interactive=True), 
            gr.update(interactive=True), 
            gr.update(value=LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"]).get("Stop Translation", "Stop Translation"), interactive=False)
        ),
        inputs=[session_lang],
        outputs=[translate_button, continue_button, stop_button]
    )

    # Update stop button handler to pass session_lang:
    stop_button.click(
        request_stop_translation,
        inputs=[session_lang],
        outputs=[stop_button]
    )

    # Language swap functionality
    def swap_languages(src_lang, dst_lang):
        """Swap source and target languages"""        
        # Update preferences with swapped values
        update_language_preferences(src_lang=dst_lang, dst_lang=src_lang)
        
        # Return swapped values
        return dst_lang, src_lang
    
    # Language change event handlers - show/hide custom language row
    src_lang.change(on_src_language_change, inputs=src_lang, outputs=[custom_lang_row])
    dst_lang.change(on_dst_language_change, inputs=dst_lang, outputs=[custom_lang_row])
    swap_button.click(swap_languages, inputs=[src_lang, dst_lang], outputs=[src_lang, dst_lang])

    # Create new language
    def on_add_new(lang_name, session_lang):
        success, msg = add_custom_language(lang_name)
        # Get translated "Add Custom Language" label for current session
        labels = LABEL_TRANSLATIONS.get(session_lang, LABEL_TRANSLATIONS["en"])
        custom_label = labels.get("Add Custom Language", "+ Add Custom...")
        new_choices = get_available_languages() + [custom_label]
        # Pick newly created language as selected value
        new_val = lang_name if success else new_choices[0]
        return (
            gr.update(choices=new_choices, value=new_val),
            gr.update(choices=new_choices, value=new_val),
            gr.update(value=""),  # Clear input
            gr.update(visible=False)  # Hide custom lang row
        )

    add_lang_button.click(
        on_add_new,
        inputs=[custom_lang_input, session_lang],
        outputs=[src_lang, dst_lang, custom_lang_input, custom_lang_row]
    )

    # Translation history event handlers - Page navigation using JavaScript
    def navigate_to_history(session_lang):
        """Navigate to history page and load records"""
        html_content = load_translation_history(session_lang)
        return html_content

    # Navigate to history page
    history_nav_btn.click(
        navigate_to_history,
        inputs=[session_lang],
        outputs=[history_list]
    ).then(
        fn=None,
        inputs=None,
        outputs=None,
        js="""
        () => {
            const mainPage = document.getElementById('main-page');
            const historyPage = document.getElementById('history-page');
            if (mainPage) mainPage.style.display = 'none';
            if (historyPage) historyPage.style.display = 'block';
        }
        """
    )

    # Navigate back to main page
    history_back_btn.click(
        fn=None,
        inputs=None,
        outputs=None,
        js="""
        () => {
            const mainPage = document.getElementById('main-page');
            const historyPage = document.getElementById('history-page');
            if (mainPage) mainPage.style.display = 'block';
            if (historyPage) historyPage.style.display = 'none';
        }
        """
    )

    # Refresh history
    history_refresh_btn.click(
        load_translation_history,
        inputs=[session_lang],
        outputs=[history_list]
    )

    # Folder opening handler
    folder_open_trigger.click(
        open_folder_path,
        inputs=[folder_path_input],
        outputs=[]
    )

    theme_toggle_btn.click(
        fn=None,
        inputs=[],
        outputs=[],
        js="""
        function() {
            // Get the root element
            const root = document.querySelector('.gradio-container').closest('body') || document.body;
            
            // Check current theme - Gradio uses 'dark' class on body
            const isDark = root.classList.contains('dark');
            
            // Toggle theme class
            if (isDark) {
                root.classList.remove('dark');
                root.classList.add('light');
            } else {
                root.classList.remove('light');
                root.classList.add('dark');
            }
            
            // Update button icon
            const btn = document.querySelector('#theme-toggle-btn');
            if (btn) {
                btn.innerHTML = isDark ? '☀️' : '🌙';
            }
            
            // Also try to set the data-theme attribute for compatibility
            const gradioContainer = document.querySelector('.gradio-container');
            if (gradioContainer) {
                gradioContainer.setAttribute('data-theme', isDark ? 'light' : 'dark');
            }
            
            // Force a style recalculation
            document.body.style.display = 'none';
            document.body.offsetHeight; // trigger reflow
            document.body.style.display = '';
            
            return [];
        }
        """
    )

    # Hidden textbox to store session language for JavaScript synchronization
    session_lang_holder = gr.Textbox(value="en", visible=False, elem_id="session-lang-holder")

    # On page load, set user language and labels
    demo.load(
        fn=init_ui,
        inputs=None,
        outputs=[
            session_lang, lan_mode_state, default_online_state, max_token_state, max_retries_state,
            excel_mode_2_state, excel_bilingual_mode_state, word_bilingual_mode_state, pdf_bilingual_mode_state, thread_count_state,
            use_online_model, model_choice, glossary_choice, glossary_upload_row,
            src_lang, dst_lang, use_online_model, lan_mode_checkbox,
            model_choice, glossary_choice, max_retries_slider, thread_count_slider,
            api_key_input, remember_key_checkbox, file_input, output_file, status_message, translate_button,
            continue_button, excel_mode_checkbox, excel_bilingual_checkbox, word_bilingual_checkbox, pdf_bilingual_checkbox, stop_button,
            custom_lang_input, add_lang_button, history_nav_btn, history_back_btn, history_refresh_btn, history_title
        ],
        js="""
        () => {
            console.log('Initializing API Key features...');

            // Function to open folder - triggers backend via hidden components
            window.openFolder = function(path) {
                if (!path) {
                    console.log('No path provided');
                    return;
                }
                console.log('Opening folder:', path);

                // Find the hidden input and button
                const pathInput = document.querySelector('#folder-path-input input, #folder-path-input textarea');
                const triggerBtn = document.querySelector('#folder-open-trigger');

                if (pathInput && triggerBtn) {
                    // Set the path value
                    pathInput.value = path;
                    // Dispatch input event to update Gradio state
                    pathInput.dispatchEvent(new Event('input', { bubbles: true }));

                    // Small delay to ensure state update, then click trigger
                    setTimeout(() => {
                        triggerBtn.click();
                    }, 100);
                } else {
                    console.log('Folder open components not found');
                    // Fallback: show path in alert
                    alert('Folder path: ' + path);
                }
            };

            // Translations for tooltip and API key label (dynamically generated from languages_config.py)
            window.apiKeyTranslations = """ + generate_api_key_translations_js() + """;

            // Store current language
            window.currentApiKeyLang = 'en';

            // Update API key language function
            window.updateApiKeyLanguage = function(lang) {
                window.currentApiKeyLang = lang || 'en';
                const trans = window.apiKeyTranslations[window.currentApiKeyLang] || window.apiKeyTranslations["en"];
                const labelEl = document.getElementById('api-key-label-text');
                const titleEl = document.getElementById('tooltip-title-text');
                const contentEl = document.getElementById('tooltip-content-text');
                if (labelEl) labelEl.textContent = trans.label;
                if (titleEl) titleEl.textContent = trans.tooltipTitle;
                if (contentEl) contentEl.textContent = trans.tooltipContent;
                console.log('API Key language updated to:', window.currentApiKeyLang);
            };

            // Function to add eye toggle button
            window.initApiKeyEyeToggle = function() {
                const apiKeyInput = document.querySelector('#api-key-input input[type="password"], #api-key-input input[type="text"]');
                if (!apiKeyInput) {
                    console.log('API key input not found, checking if section is visible...');
                    const section = document.getElementById('api-key-section');
                    console.log('API key section:', section, 'visible:', section ? section.offsetParent : null);
                    return false;
                }
                if (document.getElementById('api-key-toggle')) {
                    console.log('Eye toggle already exists');
                    return true;
                }

                console.log('Found API key input, adding eye toggle button...');

                // Find the immediate parent of the input and set position relative
                let wrapper = apiKeyInput.parentElement;
                wrapper.style.position = 'relative';
                wrapper.style.display = 'flex';
                wrapper.style.alignItems = 'center';

                const toggleBtn = document.createElement('button');
                toggleBtn.type = 'button';
                toggleBtn.id = 'api-key-toggle';
                toggleBtn.innerHTML = '<svg class="eye-open" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" width="18" height="18"><path d="M1 12s4-8 11-8 11 8 11 8-4 8-11 8-11-8-11-8z"></path><circle cx="12" cy="12" r="3"></circle></svg><svg class="eye-closed" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" width="18" height="18" style="display:none;"><path d="M17.94 17.94A10.07 10.07 0 0 1 12 20c-7 0-11-8-11-8a18.45 18.45 0 0 1 5.06-5.94M9.9 4.24A9.12 9.12 0 0 1 12 4c7 0 11 8 11 8a18.5 18.5 0 0 1-2.16 3.19m-6.72-1.07a3 3 0 1 1-4.24-4.24"></path><line x1="1" y1="1" x2="23" y2="23"></line></svg>';
                toggleBtn.style.cssText = 'position:absolute;right:8px;top:50%;transform:translateY(-50%);width:28px;height:28px;border:none;background:transparent;cursor:pointer;display:none;align-items:center;justify-content:center;color:#718096;border-radius:6px;transition:all 0.2s;padding:4px;z-index:100;';

                // Insert the button after the input, inside the wrapper
                wrapper.appendChild(toggleBtn);
                apiKeyInput.style.paddingRight = '40px';

                const updateToggleVisibility = () => {
                    const hasValue = apiKeyInput.value && apiKeyInput.value.length > 0;
                    toggleBtn.style.display = hasValue ? 'flex' : 'none';
                    console.log('Eye toggle visibility:', hasValue);
                };

                apiKeyInput.addEventListener('input', updateToggleVisibility);
                // Also listen for change events
                apiKeyInput.addEventListener('change', updateToggleVisibility);
                // Check initial value
                updateToggleVisibility();

                toggleBtn.addEventListener('click', function(e) {
                    e.preventDefault();
                    e.stopPropagation();
                    const isPassword = apiKeyInput.type === 'password';
                    apiKeyInput.type = isPassword ? 'text' : 'password';
                    this.querySelector('.eye-open').style.display = isPassword ? 'none' : 'block';
                    this.querySelector('.eye-closed').style.display = isPassword ? 'block' : 'none';
                    console.log('Toggled password visibility to:', apiKeyInput.type);
                });

                toggleBtn.addEventListener('mouseenter', function() {
                    this.style.background = 'rgba(232, 180, 184, 0.2)';
                    this.style.color = '#e8b4b8';
                });
                toggleBtn.addEventListener('mouseleave', function() {
                    this.style.background = 'transparent';
                    this.style.color = '#718096';
                });

                console.log('Eye toggle button added successfully to wrapper:', wrapper);
                return true;
            };

            // Setup tooltip positioning and visibility
            function setupTooltip() {
                const helpWrapper = document.getElementById('api-help-wrapper');
                const tooltip = document.getElementById('api-tooltip');
                if (!helpWrapper || !tooltip) {
                    console.log('Tooltip elements not found');
                    return;
                }
                if (helpWrapper._tooltipInitialized) return;
                helpWrapper._tooltipInitialized = true;

                const showTooltip = () => {
                    const rect = helpWrapper.getBoundingClientRect();
                    const tooltipWidth = 280;

                    // Make tooltip visible but transparent to measure its height
                    tooltip.style.visibility = 'hidden';
                    tooltip.style.opacity = '0';
                    tooltip.style.display = 'block';
                    const tooltipHeight = tooltip.offsetHeight;
                    tooltip.style.display = '';

                    // Gap between tooltip and icon (tight positioning)
                    const gap = 8;

                    // Position above the help icon, centered horizontally on the icon
                    const iconCenterX = rect.left + (rect.width / 2);
                    let left = iconCenterX - (tooltipWidth / 2);
                    let top = rect.top - tooltipHeight - gap;

                    // Calculate arrow position (should point to icon center)
                    let arrowLeft = 50; // percentage

                    // Keep tooltip within viewport
                    if (left < 10) {
                        // Adjust arrow position when tooltip shifts right
                        arrowLeft = ((iconCenterX - 10) / tooltipWidth) * 100;
                        left = 10;
                    }
                    if (left + tooltipWidth > window.innerWidth - 10) {
                        // Adjust arrow position when tooltip shifts left
                        const newLeft = window.innerWidth - tooltipWidth - 10;
                        arrowLeft = ((iconCenterX - newLeft) / tooltipWidth) * 100;
                        left = newLeft;
                    }

                    // Clamp arrow position
                    arrowLeft = Math.max(15, Math.min(85, arrowLeft));

                    // If tooltip would go above viewport, show below instead
                    if (top < 10) {
                        top = rect.bottom + gap;
                        // Flip arrow to top when showing below
                        tooltip.style.setProperty('--arrow-position', 'top');
                    } else {
                        tooltip.style.setProperty('--arrow-position', 'bottom');
                    }

                    tooltip.style.left = left + 'px';
                    tooltip.style.top = top + 'px';
                    tooltip.style.setProperty('--arrow-left', arrowLeft + '%');
                    tooltip.classList.add('visible');
                };

                const hideTooltip = () => {
                    tooltip.classList.remove('visible');
                };

                helpWrapper.addEventListener('mouseenter', showTooltip);
                helpWrapper.addEventListener('mouseleave', hideTooltip);
                console.log('Tooltip setup complete');
            }

            // Initialize
            function initAll() {
                // Detect browser language
                const browserLang = navigator.language.split('-')[0];
                console.log('Browser language:', navigator.language);

                if (navigator.language.startsWith('zh-TW') || navigator.language.startsWith('zh-Hant')) {
                    window.updateApiKeyLanguage('zh-Hant');
                } else if (navigator.language.startsWith('zh')) {
                    window.updateApiKeyLanguage('zh');
                } else if (window.apiKeyTranslations[browserLang]) {
                    window.updateApiKeyLanguage(browserLang);
                }

                window.initApiKeyEyeToggle();
                setupTooltip();
            }

            // Run after delay for Gradio to fully render
            setTimeout(initAll, 800);
            setTimeout(() => window.initApiKeyEyeToggle(), 1500);
            setTimeout(() => window.initApiKeyEyeToggle(), 3000);
        }
        """
    )

    # Separate event to update API key language after init
    def get_session_lang_for_js(lang):
        """Return session language for JavaScript"""
        return lang if lang else "en"

    session_lang.change(
        fn=get_session_lang_for_js,
        inputs=[session_lang],
        outputs=[session_lang_holder],
        js="""
        (lang) => {
            console.log('Session language changed to:', lang);
            function tryUpdateLanguage() {
                if (window.updateApiKeyLanguage) {
                    window.updateApiKeyLanguage(lang || 'en');
                    if (window.initApiKeyEyeToggle) {
                        window.initApiKeyEyeToggle();
                    }
                    return true;
                }
                return false;
            }
            if (!tryUpdateLanguage()) {
                setTimeout(tryUpdateLanguage, 300);
                setTimeout(tryUpdateLanguage, 800);
            }
            return lang;
        }
        """
    )

#-------------------------------------------------------------------------
# Application Launch
#-------------------------------------------------------------------------

if __name__ == "__main__":
    # Required for Windows multiprocessing (BabelDOC uses subprocess for PDF operations)
    import multiprocessing
    multiprocessing.freeze_support()

    # In server_mode, use PORT env var (Render assigns this) or default to 10000
    if server_mode:
        available_port = int(os.environ.get("PORT", 10000))
    else:
        available_port = find_available_port(start_port=9980)

    # Enable queue for progress tracking
    demo.queue()

    if server_mode or initial_lan_mode:
        demo.launch(server_name="0.0.0.0", server_port=available_port, share=False, inbrowser=not server_mode)
    else:
        demo.launch(server_port=available_port, share=False, inbrowser=True)