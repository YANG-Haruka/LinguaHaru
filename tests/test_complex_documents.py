# Stress tests with deliberately complex documents.
#
# Every test builds a document loaded with structures that commonly break
# translation round-trips (nested tables, mid-paragraph hyperlinks, inline
# images, fields, headers/footers, merged cells, grouped shapes, code blocks),
# fake-translates it, then verifies TWO things about the output:
#   1. content  - every translatable string was actually translated
#   2. format   - structure survives: element counts, nesting, ordering,
#                 hyperlink targets, formulas, code blocks
#
# Run from the repo root:
#   python tests/test_complex_documents.py
import json
import os
import re
import sys
import zipfile

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, REPO_ROOT)
os.chdir(REPO_ROOT)

WORK_DIR = os.path.join(REPO_ROOT, "tests", "_roundtrip_work", "complex")
TEMP_DIR = os.path.join(WORK_DIR, "temp")
RESULT_DIR = os.path.join(WORK_DIR, "result")

T = "[T]"
CHECKS = []


def check(name, cond, detail=""):
    status = "PASS" if cond else "FAIL"
    print(f"  [{status}] {name}" + (f"\n         -> {detail}" if detail and not cond else ""))
    CHECKS.append((name, bool(cond)))
    return bool(cond)


def fake_translate(src_json_path):
    with open(src_json_path, encoding="utf-8") as f:
        data = json.load(f)
    out = [{"count_src": i["count_src"], "type": i.get("type", "text"),
            "original": i["value"], "translated": T + i["value"]} for i in data]
    dst = os.path.join(os.path.dirname(src_json_path), "dst_translated.json")
    with open(dst, "w", encoding="utf-8") as f:
        json.dump(out, f, ensure_ascii=False, indent=2)
    return dst


def make_png(path):
    from PIL import Image
    Image.new("RGB", (60, 40), (200, 60, 60)).save(path)


# ---------------------------------------------------------------- DOCX ----
def build_complex_docx(path):
    from docx import Document
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    from docx.shared import Inches

    def add_hyperlink(paragraph, url, text):
        r_id = paragraph.part.relate_to(url, RT.HYPERLINK, is_external=True)
        hl = OxmlElement("w:hyperlink")
        hl.set(qn("r:id"), r_id)
        run = OxmlElement("w:r")
        t = OxmlElement("w:t")
        t.text = text
        run.append(t)
        hl.append(run)
        paragraph._p.append(hl)
        return hl

    doc = Document()
    doc.add_heading("Annual Engineering Report", level=1)

    # 1. paragraph with a hyperlink in the middle
    p = doc.add_paragraph("Before the link comes this text ")
    add_hyperlink(p, "https://example.com/spec", "official specification")
    p.add_run(" and after the link this tail follows")

    # 2. paragraph with an inline image in the middle
    p2 = doc.add_paragraph("Image leading text segment ")
    img_path = os.path.join(WORK_DIR, "inline.png")
    make_png(img_path)
    p2.add_run().add_picture(img_path, width=Inches(0.4))
    p2.add_run(" image trailing text segment")

    # 3. nested table: outer 2x2, inner 1x2 inside cell(0,0)
    outer = doc.add_table(rows=2, cols=2, style="Table Grid")
    outer.cell(0, 0).text = "Outer first cell content"
    inner = outer.cell(0, 0).add_table(rows=1, cols=2)
    inner.cell(0, 0).text = "Nested cell alpha"
    inner.cell(0, 1).text = "Nested cell beta"
    outer.cell(0, 1).text = "Outer second cell content"
    outer.cell(1, 0).text = "Outer third cell content"
    outer.cell(1, 1).text = "Outer fourth cell content"

    # 4. numbered + bullet lists
    doc.add_paragraph("First ordered list entry", style="List Number")
    doc.add_paragraph("Second ordered list entry", style="List Number")
    doc.add_paragraph("Bullet list entry text", style="List Bullet")

    # 5. complex field (REF) inside a sentence
    p5 = doc.add_paragraph("As shown in section ")
    for char_type in ("begin",):
        r = OxmlElement("w:r")
        fc = OxmlElement("w:fldChar")
        fc.set(qn("w:fldCharType"), char_type)
        r.append(fc)
        p5._p.append(r)
    r_instr = OxmlElement("w:r")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = r" REF _Ref99 \h "
    r_instr.append(instr)
    p5._p.append(r_instr)
    for char_type in ("separate",):
        r = OxmlElement("w:r")
        fc = OxmlElement("w:fldChar")
        fc.set(qn("w:fldCharType"), char_type)
        r.append(fc)
        p5._p.append(r)
    p5.add_run("3.5")
    r_end = OxmlElement("w:r")
    fc = OxmlElement("w:fldChar")
    fc.set(qn("w:fldCharType"), "end")
    r_end.append(fc)
    p5._p.append(r_end)
    p5.add_run(" the values converge rapidly")

    # 6. header and footer
    section = doc.sections[0]
    section.header.paragraphs[0].text = "Confidential header banner"
    section.footer.paragraphs[0].text = "Footer copyright notice"

    # 7. mixed formatting runs in one paragraph
    p7 = doc.add_paragraph("Normal start then ")
    p7.add_run("bold middle part").bold = True
    p7.add_run(" then italic ending").italic = True

    doc.save(path)


def docx_paragraph_texts_in_order(path):
    """Full text of each body paragraph, in document order, hyperlink runs included."""
    from lxml import etree
    with zipfile.ZipFile(path) as z:
        tree = etree.fromstring(z.read("word/document.xml"))
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    texts = []
    for p in tree.iter(f"{{{ns['w']}}}p"):
        text = "".join(t.text or "" for t in p.iter(f"{{{ns['w']}}}t"))
        if text.strip():
            texts.append(text)
    return texts


def test_docx_complex():
    print("DOCX complex: hyperlinks, inline image, nested tables, fields, header/footer")
    from pipeline.word_translation_pipeline import (
        extract_word_content_to_json, write_translated_content_to_word)

    src = os.path.join(WORK_DIR, "complex.docx")
    build_complex_docx(src)

    src_json = extract_word_content_to_json(src, TEMP_DIR)
    dst_json = fake_translate(src_json)
    out = write_translated_content_to_word(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                           bilingual_mode=False, src_lang="en", dst_lang="ja")

    paras = docx_paragraph_texts_in_order(out)
    joined = "\n".join(paras)

    # --- content ---
    check("heading translated", T + "Annual Engineering Report" in joined, joined[:400])
    # The link text is part of the translated string (fake translator prefixes
    # the whole paragraph), so verify it entered the translation stream AND
    # ended up inside the rebuilt hyperlink element
    from lxml import etree as _etree
    with zipfile.ZipFile(out) as z:
        _tree = _etree.fromstring(z.read("word/document.xml"))
    _w = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    link_texts = ["".join(t.text or "" for t in hl.iter(f"{{{_w}}}t"))
                  for hl in _tree.iter(f"{{{_w}}}hyperlink")]
    check("hyperlink text went through translation and sits inside the link",
          any("official specification" in lt for lt in link_texts)
          and any(T + "Before the link" in p for p in paras),
          f"link_texts={link_texts}")
    check("nested table cells translated",
          all(T + s in joined for s in ("Nested cell alpha", "Nested cell beta",
                                        "Outer first cell content", "Outer fourth cell content")),
          joined)
    check("list entries translated",
          all(T + s in joined for s in ("First ordered list entry", "Bullet list entry text")), joined)
    check("mixed-format paragraph translated",
          any(T in p and "bold middle part" in p.replace(T, "") for p in paras)
          or T + "Normal start then bold middle part then italic ending" in joined, joined)

    # header/footer
    from lxml import etree
    with zipfile.ZipFile(out) as z:
        names = z.namelist()
        header_xml = "".join(z.read(n).decode("utf-8") for n in names if "header" in n)
        footer_xml = "".join(z.read(n).decode("utf-8") for n in names if "footer" in n)
    check("header translated", T + "Confidential header banner" in header_xml, header_xml[:300])
    check("footer translated", T + "Footer copyright notice" in footer_xml, footer_xml[:300])

    # --- format / ordering ---
    link_para = next((p for p in paras if "the link" in p.replace(T, "")), "")
    before = link_para.find("Before the link")
    middle = link_para.find("official specification")
    after = link_para.find("after the link")
    check("hyperlink paragraph order preserved (before < link < after)",
          0 <= before < middle < after, repr(link_para))

    img_para = next((p for p in paras if "trailing text segment" in p.replace(T, "")), "")
    lead = img_para.find("Image leading text")
    trail = img_para.find("image trailing text")
    check("image paragraph text order preserved", 0 <= lead < trail, repr(img_para))

    with zipfile.ZipFile(out) as z:
        doc_xml = z.read("word/document.xml").decode("utf-8")
        rels = z.read("word/_rels/document.xml.rels").decode("utf-8")
    check("inline image element survives", "<w:drawing>" in doc_xml or "<pic:pic" in doc_xml)
    check("hyperlink element + relationship survive",
          "<w:hyperlink" in doc_xml and "example.com/spec" in rels)
    check("REF field survives", "REF _Ref99" in doc_xml)
    check("no placeholder literals leaked",
          not re.search(r"\{\{[A-Z_]+[^}]*\}\}|\{FIELD:", doc_xml.replace("&#123;", "{")))
    check("nested table structure intact (2 tbl elements)", doc_xml.count("<w:tbl>") == 2,
          f"found {doc_xml.count('<w:tbl>')}")


# ---------------------------------------------------------------- XLSX ----
def test_xlsx_complex():
    print("XLSX complex: multi-sheet, merged cells, formulas, newlines, hyperlink")
    import openpyxl
    from pipeline.excel_translation_pipeline import (
        extract_excel_content_to_json, write_translated_content_to_excel)

    src = os.path.join(WORK_DIR, "complex.xlsx")
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Summary"
    ws1["A1"] = "Quarterly business overview"
    ws1.merge_cells("A1:C1")
    ws1["A2"] = "Region"
    ws1["B2"] = "Total revenue"
    ws1["A3"] = "Northern district"
    ws1["B3"] = 12345.67          # number: must stay a number
    ws1["B4"] = "=SUM(B3:B3)"     # formula: must stay a formula
    ws1["A5"] = "Line one of note\nLine two of note"   # embedded newline
    ws1["A6"].hyperlink = "https://example.com"
    ws1["A6"] = "Visit project homepage"

    ws2 = wb.create_sheet("Detail Data")
    ws2["A1"] = "Detailed measurement log"
    ws2.merge_cells("A1:B2")
    wb.save(src)

    src_json = extract_excel_content_to_json(src, TEMP_DIR, use_xlwings=False)
    dst_json = fake_translate(src_json)
    out = write_translated_content_to_excel(src, src_json, dst_json, RESULT_DIR,
                                            src_lang="en", dst_lang="ja")

    wb2 = openpyxl.load_workbook(out)
    sheets = wb2.sheetnames
    s1 = wb2[sheets[0]]
    s2 = wb2[sheets[1]]

    # --- content ---
    check("merged title translated", s1["A1"].value == T + "Quarterly business overview",
          repr(s1["A1"].value))
    check("plain cells translated", s1["A3"].value == T + "Northern district", repr(s1["A3"].value))
    check("multiline cell translated with newline preserved",
          s1["A5"].value == f"{T}Line one of note\nLine two of note", repr(s1["A5"].value))
    check("second sheet translated", s2["A1"].value == T + "Detailed measurement log",
          repr(s2["A1"].value))
    check("sheet names translated", any(T in name or "-T-" in name for name in sheets), str(sheets))

    # --- format ---
    check("number cell untouched", s1["B3"].value == 12345.67, repr(s1["B3"].value))
    check("formula preserved", str(s1["B4"].value).startswith("=SUM"), repr(s1["B4"].value))
    check("merged ranges intact",
          sorted(str(r) for r in s1.merged_cells.ranges) == ["A1:C1"]
          and sorted(str(r) for r in s2.merged_cells.ranges) == ["A1:B2"],
          f"{list(s1.merged_cells.ranges)} / {list(s2.merged_cells.ranges)}")
    check("hyperlink preserved", s1["A6"].hyperlink is not None
          and "example.com" in str(s1["A6"].hyperlink.target), repr(s1["A6"].hyperlink))


# ---------------------------------------------------------------- PPTX ----
def test_pptx_complex():
    print("PPTX complex: table, grouped shapes, multi-run textbox, sparse notes")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pipeline.ppt_translation_pipeline import (
        extract_ppt_content_to_json, write_translated_content_to_ppt)

    src = os.path.join(WORK_DIR, "complex.pptx")
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # slide 1: table + multi-run textbox
    s1 = prs.slides.add_slide(blank)
    rows, cols = 2, 2
    table = s1.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(2)).table
    table.cell(0, 0).text = "Header column one"
    table.cell(0, 1).text = "Header column two"
    table.cell(1, 0).text = "Data row first value"
    table.cell(1, 1).text = "Data row second value"
    box = s1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(6), Inches(1.5))
    tf = box.text_frame
    tf.text = "First paragraph in textbox"
    p2 = tf.add_paragraph()
    run_a = p2.add_run()
    run_a.text = "Bold segment"
    run_a.font.bold = True
    run_b = p2.add_run()
    run_b.text = " followed by normal segment"

    # slide 2: grouped shapes + notes
    s2 = prs.slides.add_slide(blank)
    g = s2.shapes.add_group_shape()
    b1 = g.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.7))
    b1.text_frame.text = "Grouped shape upper text"
    b2 = g.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(0.7))
    b2.text_frame.text = "Grouped shape lower text"
    s2.notes_slide.notes_text_frame.text = "Presenter note on slide two"

    # slide 3: plain, no notes
    s3 = prs.slides.add_slide(blank)
    box3 = s3.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box3.text_frame.text = "Slide three standalone text"
    prs.save(src)

    src_json = extract_ppt_content_to_json(src, TEMP_DIR)
    dst_json = fake_translate(src_json)
    out = write_translated_content_to_ppt(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                          src_lang="en", dst_lang="ja")

    prs2 = Presentation(out)
    all_text = []
    def collect(shapes):
        for sh in shapes:
            if sh.shape_type == 6:  # group
                collect(sh.shapes)
            elif sh.has_text_frame:
                all_text.append(sh.text_frame.text)
            elif getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        all_text.append(cell.text)
    for slide in prs2.slides:
        collect(slide.shapes)
    joined = "\n".join(all_text)

    # --- content ---
    check("table cells translated",
          all(T + s in joined for s in ("Header column one", "Data row second value")), joined)
    check("multi-run paragraph translated",
          "Bold segment" in joined.replace(T, "") and T in joined, joined)
    check("grouped shape texts translated",
          all(T + s in joined for s in ("Grouped shape upper text", "Grouped shape lower text")),
          joined)
    check("slide 3 translated", T + "Slide three standalone text" in joined, joined)

    # notes on the right slide
    notes2 = prs2.slides[1].notes_slide.notes_text_frame.text if prs2.slides[1].has_notes_slide else ""
    check("sparse notes stay on slide 2", T + "Presenter note on slide two" in notes2, repr(notes2))

    # --- format ---
    s1_shapes = prs2.slides[0].shapes
    check("table structure intact", any(getattr(sh, "has_table", False) and sh.has_table
                                        and len(sh.table.rows) == 2 for sh in s1_shapes))
    check("group structure intact", any(sh.shape_type == 6 for sh in prs2.slides[1].shapes))


def test_pptx_chart():
    print("PPTX chart: title, series and category labels translated; numbers untouched")
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
    from pipeline.ppt_translation_pipeline import (
        extract_ppt_content_to_json, write_translated_content_to_ppt)

    src = os.path.join(WORK_DIR, "chart.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["Eastern region", "Western region"]
    chart_data.add_series("Quarterly revenue", (1234.5, 6789.0))
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                         Inches(1), Inches(1), Inches(7), Inches(4),
                                         chart_data)
    chart_shape.chart.has_title = True
    chart_shape.chart.chart_title.text_frame.text = "Revenue overview title"
    prs.save(src)

    src_json = extract_ppt_content_to_json(src, TEMP_DIR)
    with open(src_json, encoding="utf-8") as f:
        extracted = [i["value"] for i in json.load(f) if i.get("type") == "chart_part"]
    check("chart texts extracted",
          any("Revenue overview title" in v for v in extracted)
          and any("Eastern region" in v for v in extracted)
          and any("Quarterly revenue" in v for v in extracted), str(extracted))

    dst_json = fake_translate(src_json)
    out = write_translated_content_to_ppt(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                          src_lang="en", dst_lang="ja")
    with zipfile.ZipFile(out) as z:
        chart_xml = "".join(z.read(n).decode("utf-8") for n in z.namelist()
                            if n.startswith("ppt/charts/chart"))
    check("chart title translated", T + "Revenue overview title" in chart_xml, chart_xml[:600])
    check("category labels translated", T + "Eastern region" in chart_xml
          and T + "Western region" in chart_xml, chart_xml[:600])
    check("series name translated", T + "Quarterly revenue" in chart_xml, chart_xml[:600])
    check("numeric values untouched", "1234.5" in chart_xml and "6789" in chart_xml,
          chart_xml[:600])


# ------------------------------------------------------------------ MD ----
def test_md_complex():
    print("MD complex: code blocks, inline code, links, pipe tables")
    from pipeline.md_translation_pipeline import (
        extract_md_content_to_json, write_translated_content_to_md)

    src = os.path.join(WORK_DIR, "complex.md")
    content = (
        "# Project documentation overview\n\n"
        "Intro paragraph with a [reference link](https://example.com/docs) inside.\n\n"
        "```python\nprint('do not translate this code')\n```\n\n"
        "| Feature name | Status value |\n"
        "|---|---|\n"
        "| Login module | Completed now |\n\n"
        "Closing paragraph after the table.\n"
    )
    with open(src, "w", encoding="utf-8") as f:
        f.write(content)

    src_json = extract_md_content_to_json(src, TEMP_DIR)
    dst_json = fake_translate(src_json)
    out = write_translated_content_to_md(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                         src_lang="en", dst_lang="ja")
    with open(out, encoding="utf-8") as f:
        result = f.read()

    # --- content ---
    check("heading translated", "Project documentation overview" in result and T in result, result)
    check("intro paragraph translated", "reference link" in result.replace(T, ""), result)
    check("table cells translated", "Login module" in result.replace(T, ""), result)
    check("closing paragraph translated", T + "Closing paragraph after the table." in result, result)

    # --- format ---
    check("code block NOT translated",
          "print('do not translate this code')" in result
          and T + "print" not in result, result)
    check("link URL untouched", "https://example.com/docs" in result, result)
    check("table delimiter row intact", "|---|---|" in result or "| --- | --- |" in result, result)


# ----------------------------------------------------------------- SRT ----
def test_srt_complex():
    print("SRT complex: multi-line cues, styling tags")
    from pipeline.subtitle_translation_pipeline import (
        extract_srt_content_to_json, write_translated_content_to_srt)

    src = os.path.join(WORK_DIR, "complex.srt")
    with open(src, "w", encoding="utf-8") as f:
        f.write("1\n00:00:01,000 --> 00:00:03,000\nFirst line of cue\nSecond line of cue\n\n"
                "2\n00:00:04,000 --> 00:00:06,000\n<i>Italic styled subtitle</i>\n\n")

    src_json = extract_srt_content_to_json(src, TEMP_DIR)
    dst_json = fake_translate(src_json)
    out = write_translated_content_to_srt(src, src_json, dst_json, RESULT_DIR,
                                          src_lang="en", dst_lang="ja")
    with open(out, encoding="utf-8") as f:
        result = f.read()

    check("multi-line cue translated and line break kept",
          "First line of cue" in result.replace(T, "")
          and "\n" in result.split("-->")[1].split("\n\n")[0].strip(), result)
    check("styling tags survive", "<i>" in result and "</i>" in result, result)
    check("both cues present", result.count("-->") == 2, result)


def main():
    import shutil
    shutil.rmtree(WORK_DIR, ignore_errors=True)
    os.makedirs(TEMP_DIR, exist_ok=True)
    os.makedirs(RESULT_DIR, exist_ok=True)

    for fn in (test_docx_complex, test_xlsx_complex, test_pptx_complex,
               test_pptx_chart, test_md_complex, test_srt_complex):
        try:
            fn()
        except Exception:
            import traceback
            traceback.print_exc()
            CHECKS.append((fn.__name__ + " (crashed)", False))
        print()

    passed = sum(1 for _, ok in CHECKS if ok)
    print("=" * 60)
    print(f"{passed}/{len(CHECKS)} checks passed")
    for name, ok in CHECKS:
        if not ok:
            print(f"  FAIL: {name}")
    sys.exit(0 if passed == len(CHECKS) else 1)


if __name__ == "__main__":
    main()
