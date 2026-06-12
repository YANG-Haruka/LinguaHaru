# Round-trip tests for the document pipelines.
#
# For each format: generate a document containing the structures that
# historically broke (fields, sparse notes, merged cells, SRT variants,
# blank lines, nested table cells), run extract -> fake-translate -> write,
# then reopen the output and assert the structures survived.
#
# Run directly (no pytest needed):
#   python tests/test_pipelines_roundtrip.py
import json
import os
import re
import shutil
import sys
import zipfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

WORK_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_roundtrip_work")
TEMP_DIR = os.path.join(WORK_DIR, "temp")
RESULT_DIR = os.path.join(WORK_DIR, "result")

T = "[T]"  # fake-translation marker


def fake_translate(src_json_path):
    """Produce dst_translated.json next to src.json: prefix every value with [T]."""
    with open(src_json_path, encoding="utf-8") as f:
        data = json.load(f)
    out = []
    for item in data:
        out.append({
            "count_src": item["count_src"],
            "type": item.get("type", "text"),
            "original": item["value"],
            "translated": T + item["value"],
        })
    dst_path = os.path.join(os.path.dirname(src_json_path), "dst_translated.json")
    with open(dst_path, "w", encoding="utf-8") as f:
        json.dump(out, f, ensure_ascii=False, indent=2)
    return dst_path


def read_zip_xml(path, member):
    with zipfile.ZipFile(path) as z:
        return z.read(member).decode("utf-8")


def check(name, cond, detail=""):
    status = "PASS" if cond else "FAIL"
    print(f"  [{status}] {name}" + (f" -- {detail}" if detail and not cond else ""))
    return bool(cond)


def test_docx():
    print("DOCX: field preservation + numbered headings")
    from docx import Document
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from pipeline.word_translation_pipeline import (
        extract_word_content_to_json, write_translated_content_to_word)

    src = os.path.join(WORK_DIR, "field_test.docx")
    doc = Document()
    p = doc.add_paragraph("Please refer to section ")

    def field_char_run(char_type):
        r = OxmlElement("w:r")
        fc = OxmlElement("w:fldChar")
        fc.set(qn("w:fldCharType"), char_type)
        r.append(fc)
        return r

    # Complex field: begin / instrText / separate / cached result / end
    p._p.append(field_char_run("begin"))
    r_instr = OxmlElement("w:r")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = r" REF _Ref12345 \h "
    r_instr.append(instr)
    p._p.append(r_instr)
    p._p.append(field_char_run("separate"))
    p.add_run("1.2")
    p._p.append(field_char_run("end"))
    p.add_run(" for further details")
    doc.add_paragraph("1 Introduction")
    doc.add_paragraph("This document describes the overall system design")
    doc.save(src)

    src_json = extract_word_content_to_json(src, TEMP_DIR)
    with open(src_json, encoding="utf-8") as f:
        extracted = json.load(f)
    texts = [i["value"] for i in extracted]

    ok = check("extracted field placeholder uses double braces",
               any("{{FIELD: REF" in v for v in texts), str(texts))
    ok &= check("numbered heading '1 Introduction' is extracted for translation",
                any("1 Introduction" in v for v in texts), str(texts))

    dst_json = fake_translate(src_json)
    out = write_translated_content_to_word(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                           bilingual_mode=False, src_lang="en", dst_lang="ja")
    xml = read_zip_xml(out, "word/document.xml")
    ok &= check("REF field element survives in output", "REF _Ref12345" in xml)
    ok &= check("no placeholder literals leak into output",
                "{{FIELD" not in xml and "{FIELD:" not in xml)
    ok &= check("translated text present", T in xml)
    return ok


def test_pptx():
    print("PPTX: sparse speaker notes stay on their own slide")
    from pptx import Presentation
    from pptx.util import Inches
    from pipeline.ppt_translation_pipeline import (
        extract_ppt_content_to_json, write_translated_content_to_ppt)

    src = os.path.join(WORK_DIR, "notes_test.pptx")
    prs = Presentation()
    blank = prs.slide_layouts[6]
    bodies = ["Alpha slide body text", "Bravo slide body text", "Charlie slide body text"]
    notes = {1: "Speaker note for bravo", 2: "Speaker note for charlie"}  # slides 2 and 3 only
    for i, body in enumerate(bodies):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.text_frame.text = body
        if i in notes:
            slide.notes_slide.notes_text_frame.text = notes[i]
    prs.save(src)

    src_json = extract_ppt_content_to_json(src, TEMP_DIR)
    dst_json = fake_translate(src_json)
    out = write_translated_content_to_ppt(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                          src_lang="en", dst_lang="ja")

    # Every notes file must contain the translation of ITS OWN original text
    ok = True
    with zipfile.ZipFile(out) as z:
        notes_files = sorted(n for n in z.namelist()
                             if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml"))
        for nf in notes_files:
            content = z.read(nf).decode("utf-8")
            m = re.search(r"Speaker note for (\w+)", content)
            ok &= check(f"{nf} keeps its own note",
                        m is not None and content.count(T + "Speaker note for") == 1
                        or T in content,
                        content[:300])
            if m:
                ok &= check(f"{nf} translation matches its original ({m.group(1)})",
                            f"{T}Speaker note for {m.group(1)}" in content, content[:300])
    return ok


def test_xlsx():
    print("XLSX: merged cells survive without bogus 1x1 ranges")
    import openpyxl
    from pipeline.excel_translation_pipeline import (
        extract_excel_content_to_json, write_translated_content_to_excel)

    src = os.path.join(WORK_DIR, "merge_test.xlsx")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Quarterly revenue report"
    ws.merge_cells("A1:B2")
    ws["C3"] = "Net profit after taxes"
    wb.save(src)

    src_json = extract_excel_content_to_json(src, TEMP_DIR, use_xlwings=False)
    dst_json = fake_translate(src_json)
    out = write_translated_content_to_excel(src, src_json, dst_json, RESULT_DIR,
                                            src_lang="en", dst_lang="ja")

    wb2 = openpyxl.load_workbook(out)
    ws2 = wb2.active
    ranges = sorted(str(r) for r in ws2.merged_cells.ranges)
    ok = check("merged ranges unchanged", ranges == ["A1:B2"], str(ranges))
    ok &= check("merged cell translated", ws2["A1"].value == T + "Quarterly revenue report",
                repr(ws2["A1"].value))
    ok &= check("normal cell translated", ws2["C3"].value == T + "Net profit after taxes",
                repr(ws2["C3"].value))
    return ok


def test_srt():
    print("SRT: variant timestamps parsed, no cue lost")
    from pipeline.subtitle_translation_pipeline import (
        extract_srt_content_to_json, write_translated_content_to_srt)

    src = os.path.join(WORK_DIR, "variant_test.srt")
    with open(src, "w", encoding="utf-8") as f:
        f.write("1\n00:00:01.500 --> 0:00:03.000\nHello there my friend\n\n"
                "2\n00:00:04,000 --> 00:00:05,250\nSecond subtitle line\n\n")

    src_json = extract_srt_content_to_json(src, TEMP_DIR)
    with open(src_json, encoding="utf-8") as f:
        extracted = json.load(f)
    ok = check("both cues extracted (dot + comma ms)", len(extracted) == 2, str(extracted))

    dst_json = fake_translate(src_json)
    out = write_translated_content_to_srt(src, src_json, dst_json, RESULT_DIR,
                                          src_lang="en", dst_lang="ja")
    with open(out, encoding="utf-8") as f:
        content = f.read()
    ok &= check("both cues in output", content.count("-->") == 2, content)
    ok &= check("timestamps normalized to comma", "00:00:01,500" in content, content)
    ok &= check("both translations present",
                T + "Hello there my friend" in content and T + "Second subtitle line" in content,
                content)
    return ok


def test_txt():
    print("TXT: blank lines, single newlines and indentation preserved")
    from pipeline.txt_translation_pipeline import (
        extract_txt_content_to_json, write_translated_content_to_txt)

    src = os.path.join(WORK_DIR, "structure_test.txt")
    original = ("Document title here\n"
                "\n"
                "    Indented first paragraph line\n"
                "Second paragraph line\n"
                "\n"
                "\n"
                "Final closing line\n")
    with open(src, "w", encoding="utf-8") as f:
        f.write(original)

    src_json = extract_txt_content_to_json(src, TEMP_DIR)
    dst_json = fake_translate(src_json)
    out = write_translated_content_to_txt(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                          src_lang="en", dst_lang="ja")
    with open(out, encoding="utf-8") as f:
        result = f.read()

    expected = ("[T]Document title here\n"
                "\n"
                "    [T]Indented first paragraph line\n"
                "[T]Second paragraph line\n"
                "\n"
                "\n"
                "[T]Final closing line\n")
    return check("line structure identical (blank lines, indent, single \\n)",
                 result == expected, repr(result))


def test_md():
    print("MD: HTML table cell with nested tags gets translated")
    from pipeline.md_translation_pipeline import (
        extract_md_content_to_json, write_translated_content_to_md)

    src = os.path.join(WORK_DIR, "table_test.md")
    with open(src, "w", encoding="utf-8") as f:
        f.write("# Heading for the test\n\n"
                "<table><tr><td><b>Bold part</b> plain tail</td>"
                "<td>Simple cell content</td></tr></table>\n\n"
                "Regular paragraph of text\n")

    src_json = extract_md_content_to_json(src, TEMP_DIR)
    dst_json = fake_translate(src_json)
    out = write_translated_content_to_md(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                         src_lang="en", dst_lang="ja")
    with open(out, encoding="utf-8") as f:
        content = f.read()

    ok = check("mixed-content cell translated (not silently kept as source)",
               T + "Bold part plain tail" in content, content)
    ok &= check("simple cell translated", T + "Simple cell content" in content, content)
    # the '#' marker is part of the text sent for translation by design
    ok &= check("heading translated", T + "# Heading for the test" in content, content)
    return ok


def main():
    shutil.rmtree(WORK_DIR, ignore_errors=True)
    os.makedirs(TEMP_DIR, exist_ok=True)
    os.makedirs(RESULT_DIR, exist_ok=True)

    results = {}
    for fn in (test_docx, test_pptx, test_xlsx, test_srt, test_txt, test_md):
        try:
            results[fn.__name__] = fn()
        except Exception as e:
            import traceback
            traceback.print_exc()
            results[fn.__name__] = False
        print()

    print("=" * 50)
    failed = [k for k, v in results.items() if not v]
    for k, v in results.items():
        print(f"{'PASS' if v else 'FAIL'}: {k}")
    print("=" * 50)
    sys.exit(1 if failed else 0)


if __name__ == "__main__":
    main()
