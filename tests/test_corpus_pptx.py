# Corpus tests: PPTX structures.
#   - table with merged cells (gridSpan / vMerge)
#   - group shape nested inside another group shape
#   - chart with two series
#   - multi-paragraph speaker notes
#   - bullet indent levels
#
# Run from the repo root:
#   python tests/test_corpus_pptx.py
import os
import sys
import zipfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from tests.corpus_common import T, check, fake_translate, run, work_dirs

WORK_DIR, TEMP_DIR, RESULT_DIR = work_dirs("pptx")


def build_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    blank = prs.slide_layouts[6]

    # --- slide 1: table with merged cells ---
    slide1 = prs.slides.add_slide(blank)
    table = slide1.shapes.add_table(3, 3, Inches(1), Inches(1),
                                    Inches(6), Inches(2)).table
    table.cell(0, 0).merge(table.cell(0, 2))          # gridSpan across row 0
    table.cell(1, 0).merge(table.cell(2, 0))          # vMerge down column 0
    table.cell(0, 0).text = "Merged header across columns"
    table.cell(1, 0).text = "Merged label down rows"
    table.cell(1, 1).text = "Body cell one"
    table.cell(1, 2).text = "Body cell two"
    table.cell(2, 1).text = "Body cell three"
    table.cell(2, 2).text = "Body cell four"

    # --- slide 2: group inside group ---
    slide2 = prs.slides.add_slide(blank)
    outer = slide2.shapes.add_group_shape()
    box_a = outer.shapes.add_textbox(Emu(914400), Emu(914400),
                                     Emu(2743200), Emu(457200))
    box_a.text_frame.text = "Outer group caption text"
    inner = outer.shapes.add_group_shape()
    box_b = inner.shapes.add_textbox(Emu(914400), Emu(1828800),
                                     Emu(2743200), Emu(457200))
    box_b.text_frame.text = "Inner group caption text"

    # --- slide 3: chart with two series + multi-paragraph notes ---
    slide3 = prs.slides.add_slide(blank)
    chart_data = CategoryChartData()
    chart_data.categories = ["Spring quarter", "Summer quarter"]
    chart_data.add_series("Revenue series total", (110.0, 230.0))
    chart_data.add_series("Expense series total", (90.0, 150.0))
    gframe = slide3.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                     Inches(1), Inches(1), Inches(7), Inches(4),
                                     chart_data)
    gframe.chart.has_title = True
    gframe.chart.chart_title.text_frame.text = "Quarterly outcome chart"
    notes = slide3.notes_slide.notes_text_frame
    notes.text = "First paragraph of the speaker notes"
    notes.add_paragraph().text = "Second paragraph with extra reminders"
    notes.add_paragraph().text = "Third paragraph closing the notes"

    # --- slide 4: bullet indent levels ---
    slide4 = prs.slides.add_slide(blank)
    box = slide4.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
    tf = box.text_frame
    tf.text = "Top level bullet point"
    p1 = tf.add_paragraph()
    p1.text = "Second level bullet point"
    p1.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Third level bullet point"
    p2.level = 2

    prs.save(path)


def test_pptx_structures():
    print("PPTX: merged table, nested groups, 2-series chart, notes, bullet levels")
    from pptx import Presentation
    from core.pipelines.ppt_translation_pipeline import (
        extract_ppt_content_to_json, write_translated_content_to_ppt)

    src = os.path.join(WORK_DIR, "structures.pptx")
    build_pptx(src)

    src_json = extract_ppt_content_to_json(src, TEMP_DIR)
    dst_json = fake_translate(src_json)
    out = write_translated_content_to_ppt(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                          src_lang="en", dst_lang="ja")

    prs = Presentation(out)
    slides = list(prs.slides)

    # --- slide 1: merged table ---
    table = next(s for s in slides[0].shapes if s.has_table).table
    check("gridSpan merge survives (origin + spanned)",
          table.cell(0, 0).is_merge_origin and table.cell(0, 2).is_spanned,
          f"origin={table.cell(0,0).is_merge_origin} spanned={table.cell(0,2).is_spanned}")
    check("vMerge survives", table.cell(1, 0).is_merge_origin and table.cell(2, 0).is_spanned,
          f"origin={table.cell(1,0).is_merge_origin} spanned={table.cell(2,0).is_spanned}")
    check("merged cells translated",
          table.cell(0, 0).text == T + "Merged header across columns"
          and table.cell(1, 0).text == T + "Merged label down rows",
          f"{table.cell(0,0).text!r} / {table.cell(1,0).text!r}")
    check("plain cells translated",
          table.cell(1, 1).text == T + "Body cell one"
          and table.cell(2, 2).text == T + "Body cell four",
          f"{table.cell(1,1).text!r} / {table.cell(2,2).text!r}")

    # --- slide 2: nested groups ---
    with zipfile.ZipFile(out) as z:
        slide2_xml = z.read("ppt/slides/slide2.xml").decode("utf-8")
    check("nested group structure survives (2 grpSp elements)",
          slide2_xml.count("<p:grpSp>") == 2, slide2_xml[:500])
    check("outer group textbox translated", T + "Outer group caption text" in slide2_xml,
          slide2_xml)
    check("inner (nested) group textbox translated",
          T + "Inner group caption text" in slide2_xml, slide2_xml)

    # --- slide 3: chart ---
    chart = next(s for s in slides[2].shapes if s.has_chart).chart
    check("chart title translated",
          chart.chart_title.text_frame.text == T + "Quarterly outcome chart",
          repr(chart.chart_title.text_frame.text))
    series_names = [s.name for s in chart.plots[0].series]
    check("both series names translated",
          series_names == [T + "Revenue series total", T + "Expense series total"],
          str(series_names))
    cats = list(chart.plots[0].categories)
    check("category labels translated",
          cats == [T + "Spring quarter", T + "Summer quarter"], str(cats))
    values = [list(s.values) for s in chart.plots[0].series]
    check("series values untouched",
          values == [[110.0, 230.0], [90.0, 150.0]], str(values))

    # --- slide 3 notes: every paragraph translated, order kept ---
    notes_text = slides[2].notes_slide.notes_text_frame.text
    check("all three notes paragraphs translated in order",
          notes_text.split("\n") == [T + "First paragraph of the speaker notes",
                                     T + "Second paragraph with extra reminders",
                                     T + "Third paragraph closing the notes"],
          repr(notes_text))

    # --- slide 4: bullet levels ---
    box = next(s for s in slides[3].shapes if s.has_text_frame
               and "bullet point" in s.text_frame.text.replace(T, ""))
    paras = box.text_frame.paragraphs
    check("bullet texts translated",
          [p.text for p in paras] == [T + "Top level bullet point",
                                      T + "Second level bullet point",
                                      T + "Third level bullet point"],
          str([p.text for p in paras]))
    check("bullet indent levels preserved",
          [p.level for p in paras] == [0, 1, 2], str([p.level for p in paras]))


_STATIC_SP = (
    '<p:sp><p:nvSpPr><p:cNvPr id="987" name="StaticBox"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
    '<p:spPr/><p:txBody><a:bodyPr/><a:p><a:r><a:t>Static footer text in master</a:t>'
    "</a:r></a:p></p:txBody></p:sp>")


def _inject_master_static_text(pptx_path):
    """Add a non-placeholder static text box to slideMaster1.xml."""
    import shutil
    tmp = pptx_path + ".inj"
    with zipfile.ZipFile(pptx_path) as zin, zipfile.ZipFile(tmp, "w") as zout:
        for n in zin.namelist():
            data = zin.read(n)
            if n == "ppt/slideMasters/slideMaster1.xml":
                data = data.replace(b"</p:spTree>", _STATIC_SP.encode("utf-8") + b"</p:spTree>")
            zout.writestr(n, data)
    shutil.move(tmp, pptx_path)


def test_pptx_master_static_text():
    print("PPTX: static (non-placeholder) master text translated; prompts left alone")
    from core.pipelines.ppt_translation_pipeline import (
        extract_ppt_content_to_json, write_translated_content_to_ppt)

    src = os.path.join(WORK_DIR, "master.pptx")
    build_pptx(src)
    _inject_master_static_text(src)

    src_json = extract_ppt_content_to_json(src, TEMP_DIR)
    import json
    with open(src_json, encoding="utf-8") as f:
        extracted = [i["value"] for i in json.load(f)]
    check("static master text extracted",
          "Static footer text in master" in extracted, str(extracted))
    check("placeholder prompt text NOT extracted",
          not any("Click to edit" in v for v in extracted), str(extracted))

    dst_json = fake_translate(src_json)
    out = write_translated_content_to_ppt(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                          src_lang="en", dst_lang="ja")
    with zipfile.ZipFile(out) as z:
        master = z.read("ppt/slideMasters/slideMaster1.xml").decode("utf-8")
    check("static master text translated in output",
          T + "Static footer text in master" in master, master)


# --- ADDITIVE feature fixtures: alt text, comments, notes master -----------

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

# Legacy comment part (ppt/comments/comment1.xml): body lives in <p:text>.
_LEGACY_COMMENT_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    f'<p:cmLst xmlns:p="{_P_NS}" xmlns:a="{_A_NS}">'
    '<p:cm authorId="0" dt="2024-01-01T00:00:00" idx="1">'
    '<p:pos x="100" y="100"/>'
    '<p:text>Reviewer comment needs translation</p:text>'
    "</p:cm></p:cmLst>")

# notesMaster body text box (non-placeholder static text).
_NOTES_MASTER_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    f'<p:notesMaster xmlns:p="{_P_NS}" xmlns:a="{_A_NS}"><p:cSld><p:spTree>'
    '<p:sp><p:nvSpPr><p:cNvPr id="11" name="NotesStatic"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
    '<p:spPr/><p:txBody><a:bodyPr/><a:p><a:r><a:t>Notes master footer line</a:t>'
    "</a:r></a:p></p:txBody></p:sp>"
    "</p:spTree></p:cSld></p:notesMaster>")


def _set_shape_alt_text(pptx_path):
    """Add a textbox with cNvPr descr/title alt text to slide 1 via python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation(pptx_path)
    slide1 = list(prs.slides)[0]
    box = slide1.shapes.add_textbox(Inches(5), Inches(5), Inches(2), Inches(1))
    box.text_frame.text = "Plain run text"
    cnv = box._element.find(f".//{{{_P_NS}}}cNvPr")
    cnv.set("descr", "Accessibility description of the picture")
    cnv.set("title", "Picture alt title")
    prs.save(pptx_path)


def _inject_comments_and_notesmaster(pptx_path):
    """Inject a legacy comments part and a notesMaster part via raw zip edits."""
    import shutil
    tmp = pptx_path + ".inj2"
    comment_part = "ppt/comments/comment1.xml"
    notes_master_part = "ppt/notesMasters/notesMaster1.xml"
    with zipfile.ZipFile(pptx_path) as zin, zipfile.ZipFile(tmp, "w") as zout:
        for n in zin.namelist():
            data = zin.read(n)
            if n == "[Content_Types].xml":
                # Register the two new parts so PowerPoint considers them valid.
                overrides = (
                    f'<Override PartName="/{comment_part}" ContentType='
                    '"application/vnd.openxmlformats-officedocument.presentationml.comments+xml"/>'
                    f'<Override PartName="/{notes_master_part}" ContentType='
                    '"application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"/>'
                ).encode("utf-8")
                data = data.replace(b"</Types>", overrides + b"</Types>")
            zout.writestr(n, data)
        zout.writestr(comment_part, _LEGACY_COMMENT_XML)
        zout.writestr(notes_master_part, _NOTES_MASTER_XML)
    shutil.move(tmp, pptx_path)


def test_pptx_additive_coverage():
    print("PPTX additive: shape alt text, comments, notes master")
    import json
    from core.pipelines.ppt_translation_pipeline import (
        extract_ppt_content_to_json, write_translated_content_to_ppt)

    src = os.path.join(WORK_DIR, "additive.pptx")
    build_pptx(src)
    _set_shape_alt_text(src)
    _inject_comments_and_notesmaster(src)

    src_json = extract_ppt_content_to_json(src, TEMP_DIR)
    with open(src_json, encoding="utf-8") as f:
        items = json.load(f)
    by_type = {}
    for i in items:
        by_type.setdefault(i["type"], []).append(i["value"])

    check("alt-text descr extracted",
          "Accessibility description of the picture" in by_type.get("ppt_alttext", []),
          str(by_type.get("ppt_alttext")))
    check("alt-text title extracted",
          "Picture alt title" in by_type.get("ppt_alttext", []),
          str(by_type.get("ppt_alttext")))
    check("comment text extracted",
          "Reviewer comment needs translation" in by_type.get("ppt_comment", []),
          str(by_type.get("ppt_comment")))
    check("notes master text extracted",
          "Notes master footer line" in by_type.get("ppt_notesmaster", []),
          str(by_type.get("ppt_notesmaster")))

    dst_json = fake_translate(src_json)
    out = write_translated_content_to_ppt(src, src_json, dst_json, TEMP_DIR, RESULT_DIR,
                                          src_lang="en", dst_lang="ja")

    with zipfile.ZipFile(out) as z:
        slide1 = z.read("ppt/slides/slide1.xml").decode("utf-8")
        comment = z.read("ppt/comments/comment1.xml").decode("utf-8")
        notes_master = z.read("ppt/notesMasters/notesMaster1.xml").decode("utf-8")

    check("alt-text descr translated in output",
          f'descr="{T}Accessibility description of the picture"' in slide1, slide1[:1500])
    check("alt-text title translated in output",
          f'title="{T}Picture alt title"' in slide1, slide1[:1500])
    check("comment translated in output",
          T + "Reviewer comment needs translation" in comment, comment)
    check("notes master text translated in output",
          T + "Notes master footer line" in notes_master, notes_master)

    # Regression: existing slide text still translated, structure intact.
    from pptx import Presentation
    prs = Presentation(out)
    table = next(s for s in list(prs.slides)[0].shapes if s.has_table).table
    check("existing table text still translated (no regression)",
          table.cell(0, 0).text == T + "Merged header across columns",
          repr(table.cell(0, 0).text))


if __name__ == "__main__":
    run([test_pptx_structures, test_pptx_master_static_text, test_pptx_additive_coverage])
